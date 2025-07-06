# regionalgtm_processor.py
# This module wraps the existing web_app.py logic for the Regional GTM workflow
# WITHOUT changing any of the core AI/slide selection/formatting logic

import os
import io
import json
import tempfile
from pptx import Presentation
import requests
from datetime import datetime
import traceback

# Import all the existing functions from web_app (without changing them)
from web_app import (
    get_all_slide_data,
    find_slide_by_ai, 
    analyze_and_map_content,
    populate_slide,
    analyze_template_formatting,
    deep_copy_slide_content
)

class RegionalGTMProcessor:
    """
    Orchestrates the Regional GTM workflow using existing web_app.py logic
    """
    
    def __init__(self, api_key):
        self.api_key = api_key
        self.session_data = {}
        
    def analyze_output_structure(self, template_files_data, gtm_file_data, gtm_template_file_data):
        """
        Step 1: Analyze files and present the output structure before processing
        """
        print("[REGIONAL_GTM] Analyzing output structure...")
        
        structure_info = {
            "modules": [
                {"name": "Objectives", "description": "Extract objectives from GTM template"},
                {"name": "Insights", "description": "AI-generated market research + custom insights"},
                {"name": "Timeline", "description": "Regional timeline and activations"},
                {"name": "Activation", "description": "Three-column activation plan"}
            ],
            "template_slides_count": 0,
            "gtm_slides_count": 0,
            "gtm_template_slides_count": 0
        }
        
        try:
            # Analyze template files
            for file_data in template_files_data:
                file_bytes = file_data['content']
                file_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                slides_data = get_all_slide_data(file_bytes, file_type)
                structure_info["template_slides_count"] += len(slides_data)
            
            # Analyze GTM file
            gtm_slides_data = get_all_slide_data(gtm_file_data['content'], 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
            structure_info["gtm_slides_count"] = len(gtm_slides_data)
            
            # Analyze GTM template file
            gtm_template_slides_data = get_all_slide_data(gtm_template_file_data['content'], 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
            structure_info["gtm_template_slides_count"] = len(gtm_template_slides_data)
            
            print(f"[REGIONAL_GTM] Structure analysis complete: {structure_info}")
            return structure_info
            
        except Exception as e:
            print(f"[REGIONAL_GTM] Error analyzing structure: {e}")
            return structure_info
    
    def process_objectives_module(self, template_files_data, gtm_template_file_data):
        """
        Process the Objectives module using existing AI logic
        """
        print("[REGIONAL_GTM] Processing Objectives module...")
        
        try:
            # Use existing logic to find objective slides in GTM template
            result = find_slide_by_ai(
                self.api_key,
                gtm_template_file_data['content'],
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                "objectives",
                "GTM Template (Objectives Source)"
            )
            
            objectives_content = {
                "slide_found": result["slide"] is not None,
                "slide_index": result.get("index", -1),
                "extracted_content": result["slide"] if result["slide"] else None,
                "ai_reasoning": result.get("justification", "")
            }
            
            print(f"[REGIONAL_GTM] Objectives module processed: {objectives_content['slide_found']}")
            return objectives_content
            
        except Exception as e:
            print(f"[REGIONAL_GTM] Error processing objectives: {e}")
            return {"slide_found": False, "error": str(e)}
    
    def process_insights_module(self, region, custom_insights=""):
        """
        Process the Insights module with AI-generated market research
        """
        print(f"[REGIONAL_GTM] Processing Insights module for {region}...")
        
        try:
            # AI-powered market research (similar to regionalgtm.py)
            market_research = self._get_market_research(region)
            
            insights_content = {
                "region": region,
                "ai_research": market_research,
                "custom_insights": custom_insights,
                "summary": self._summarize_insights(market_research, custom_insights)
            }
            
            print(f"[REGIONAL_GTM] Insights module processed for {region}")
            return insights_content
            
        except Exception as e:
            print(f"[REGIONAL_GTM] Error processing insights: {e}")
            return {"error": str(e)}
    
    def process_timeline_module(self, timeline_events):
        """
        Process the Timeline module
        """
        print("[REGIONAL_GTM] Processing Timeline module...")
        
        timeline_content = {
            "events": timeline_events,
            "formatted_timeline": "\n".join(f"• {event}" for event in timeline_events)
        }
        
        print(f"[REGIONAL_GTM] Timeline module processed with {len(timeline_events)} events")
        return timeline_content
    
    def process_activation_module(self, insights_summary, activation_plan, measurement_plan):
        """
        Process the Activation module
        """
        print("[REGIONAL_GTM] Processing Activation module...")
        
        activation_content = {
            "insights_summary": insights_summary,
            "activation_plan": activation_plan,
            "measurement_plan": measurement_plan
        }
        
        print("[REGIONAL_GTM] Activation module processed")
        return activation_content
    
    def assemble_final_presentation(self, template_files_data, module_results, project_name, region):
        """
        Final assembly using existing logic but only populating relevant slides
        """
        print("[REGIONAL_GTM] Assembling final presentation...")
        
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                # Load base template (use existing logic)
                new_prs = None
                all_template_slides_for_ai = []
                
                for i, file_data in enumerate(template_files_data):
                    filename = file_data['filename']
                    file_bytes = file_data['content']
                    
                    if i == 0:
                        # Use first file as base
                        new_prs = Presentation(io.BytesIO(file_bytes))
                        print(f"[REGIONAL_GTM] Loaded base template with {len(new_prs.slides)} slides")
                    else:
                        # Merge additional files (use existing logic)
                        current_prs = Presentation(io.BytesIO(file_bytes))
                        for slide_to_merge in current_prs.slides:
                            matching_layout = new_prs.slide_layouts[0]  # fallback
                            new_slide = new_prs.slides.add_slide(matching_layout)
                            deep_copy_slide_content(new_slide, slide_to_merge)
                    
                    # Get slides for AI analysis
                    slides_data = get_all_slide_data(file_bytes, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
                    all_template_slides_for_ai.extend(slides_data)
                
                # Analyze template formatting (use existing logic)
                template_formatting = analyze_template_formatting(new_prs)
                
                # Track which slides get populated
                populated_slide_indices = []
                
                # Process each module and populate slides
                modules_to_process = [
                    ("Objectives", module_results.get('objectives')),
                    ("Insights", module_results.get('insights')),
                    ("Timeline", module_results.get('timeline')),
                    ("Activation", module_results.get('activation'))
                ]
                
                for module_name, module_data in modules_to_process:
                    if not module_data:
                        continue
                        
                    print(f"[REGIONAL_GTM] Processing {module_name} module for slide population...")
                    
                    # Use existing AI logic to select best template slide
                    content_for_ai = self._prepare_content_for_ai(module_name, module_data)
                    
                    ai_mapping_result = analyze_and_map_content(
                        self.api_key,
                        content_for_ai,
                        all_template_slides_for_ai,
                        module_name.lower()
                    )
                    
                    selected_slide_index = ai_mapping_result["best_template_index"] - 1  # Convert to 0-based
                    
                    if 0 <= selected_slide_index < len(new_prs.slides):
                        dest_slide = new_prs.slides[selected_slide_index]
                        
                        # Use existing populate_slide logic with formatting
                        populate_slide(
                            dest_slide, 
                            ai_mapping_result["processed_content"],
                            template_formatting.get('title_format'),
                            template_formatting.get('body_format')
                        )
                        
                        populated_slide_indices.append(selected_slide_index)
                        print(f"[REGIONAL_GTM] Populated slide {selected_slide_index + 1} for {module_name}")
                
                # Save final presentation
                result_path = os.path.join(tmpdir, f"{project_name}_Regional_GTM.pptx")
                new_prs.save(result_path)
                
                # Read file for return
                with open(result_path, 'rb') as f:
                    pptx_bytes = f.read()
                
                print(f"[REGIONAL_GTM] Final presentation assembled with {len(populated_slide_indices)} populated slides")
                return {
                    "pptx_bytes": pptx_bytes,
                    "populated_slides": populated_slide_indices,
                    "total_slides": len(new_prs.slides)
                }
                
        except Exception as e:
            print(f"[REGIONAL_GTM] Error assembling presentation: {e}")
            raise e
    
    def _get_market_research(self, region):
        """AI-powered market research (from regionalgtm.py logic)"""
        try:
            prompt = f"You are a market research analyst. Provide a market analysis for the tech industry in {region}. Return ONLY a valid JSON object with the following keys: 'market_size', 'key_trends' (as a list of strings), 'consumer_behavior', and 'competitor_landscape' (as a list of strings)."
            
            # Use OpenAI API (similar to existing logic)
            import openai
            client = openai.OpenAI(api_key=self.api_key)
            
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                response_format={"type": "json_object"}
            )
            
            return json.loads(response.choices[0].message.content)
            
        except Exception as e:
            print(f"[REGIONAL_GTM] Error getting market research: {e}")
            return {}
    
    def _summarize_insights(self, market_research, custom_insights):
        """AI-powered insights summarization"""
        try:
            insights_text = json.dumps(market_research) + " " + custom_insights
            prompt = f"Summarize the following market research data in one or two sentences, capturing the most critical takeaway for a marketing team: '{insights_text}'"
            
            import openai
            client = openai.OpenAI(api_key=self.api_key)
            
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}]
            )
            
            return response.choices[0].message.content
            
        except Exception as e:
            print(f"[REGIONAL_GTM] Error summarizing insights: {e}")
            return "Summary could not be generated."
    
    def _prepare_content_for_ai(self, module_name, module_data):
        """Prepare module data for AI analysis"""
        if module_name == "Objectives" and module_data.get('extracted_content'):
            return {
                "title": "Objectives",
                "body": module_data['extracted_content'].get('text', ''),
                "image_data": module_data['extracted_content'].get('image_data', '')
            }
        elif module_name == "Insights":
            insights_text = f"Market Research: {json.dumps(module_data.get('ai_research', {}))}\nCustom Insights: {module_data.get('custom_insights', '')}"
            return {
                "title": f"Market Insights: {module_data.get('region', '')}",
                "body": insights_text
            }
        elif module_name == "Timeline":
            return {
                "title": "Project Timeline",
                "body": module_data.get('formatted_timeline', '')
            }
        elif module_name == "Activation":
            activation_text = f"Insights Summary: {module_data.get('insights_summary', '')}\n\nActivation Plan: {module_data.get('activation_plan', '')}\n\nMeasurement & KPIs: {module_data.get('measurement_plan', '')}"
            return {
                "title": "Regional Activation Plan",
                "body": activation_text
            }
        
        return {"title": module_name, "body": "Content not available"}
