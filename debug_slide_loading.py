#!/usr/bin/env python3
"""
Debug script to test slide loading issue with python-pptx
"""

import io
from pptx import Presentation

def test_slide_loading(template_file_path):
    """Test slide loading from a template file"""
    print(f"Testing slide loading from: {template_file_path}")
    
    # Test 1: Load directly from file path
    try:
        prs_from_path = Presentation(template_file_path)
        print(f"✓ Direct file loading: {len(prs_from_path.slides)} slides")
        for i, slide in enumerate(prs_from_path.slides):
            print(f"  Slide {i+1}: {len(slide.shapes)} shapes")
    except Exception as e:
        print(f"✗ Direct file loading failed: {e}")
        return False
    
    # Test 2: Load from bytes (like the web app does)
    try:
        with open(template_file_path, 'rb') as f:
            file_bytes = f.read()
        print(f"Read {len(file_bytes)} bytes from file")
        
        prs_from_bytes = Presentation(io.BytesIO(file_bytes))
        print(f"✓ Bytes loading: {len(prs_from_bytes.slides)} slides")
        for i, slide in enumerate(prs_from_bytes.slides):
            print(f"  Slide {i+1}: {len(slide.shapes)} shapes")
            
        # Check if both methods give same result
        if len(prs_from_path.slides) == len(prs_from_bytes.slides):
            print("✓ Both methods load the same number of slides")
            return True
        else:
            print("✗ Mismatch between file path and bytes loading!")
            return False
            
    except Exception as e:
        print(f"✗ Bytes loading failed: {e}")
        return False

if __name__ == "__main__":
    # Test with the known template file
    test_slide_loading("/workspaces/presemulator/TEST.pptx")
