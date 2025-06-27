# app.py
import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import io
import json
import requests # Used for making API calls

# --- Page Configuration ---
st.set_page_config(
    page_title="DreamAI Setups",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded",
)

# --- App Styling ---
def local_css(file_name):
    with open(file_name) as f:
        st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)

# local_css("style.css") # Placeholder for custom CSS

# --- LIVE API FUNCTIONS ---

def call_gemini_api(payload, api_key):
    """
    Generic function to call the Google Gemini API.
    """
    if not api_key:
        st.error("API Key not found. Please enter your API key in the sidebar.")
        return None

    api_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={api_key}"
    headers = {'Content-Type': 'application/json'}
    
    try:
        response = requests.post(api_url, headers=headers, json=payload)
        response.raise_for_status() # Raises an exception for bad status codes (4xx or 5xx)
        return response.json()
    except requests.exceptions.RequestException as e:
        st.error(f"API Request Failed: {e}")
        try:
            # Try to get more specific error from API response
            error_details = e.response.json()
            st.error(f"API Error Details: {error_details.get('error', {}).get('message', 'No details')}")
        except:
            pass
        return None

def get_deep_research(region, api_key):
    """
    Calls the Gemini API to get structured market research data.
    """
    st.info(f"🤖 Calling Generative AI for Deep Research on {region}...")
    
    prompt = f"Provide a market analysis for the tech industry in {region}. Give me a JSON object with the following keys: 'market_size', 'key_trends' (as a list of strings), 'consumer_behavior', and 'competitor_landscape' (as a list of strings)."
    
    payload = {
        "contents": [{"role": "user", "parts": [{"text": prompt}]}],
        "generationConfig": {
            "response_mime_type": "application/json",
            "response_schema": {
                "type": "OBJECT",
                "properties": {
                    "market_size": {"type": "STRING"},
                    "key_trends": {"type": "ARRAY", "items": {"type": "STRING"}},
                    "consumer_behavior": {"type": "STRING"},
                    "competitor_landscape": {"type": "ARRAY", "items": {"type": "STRING"}}
                }
            }
        }
    }
    
    result = call_gemini_api(payload, api_key)
    
    if result and 'candidates' in result:
        try:
            json_text = result['candidates'][0]['content']['parts'][0]['text']
            return json.loads(json_text)
        except (KeyError, IndexError, json.JSONDecodeError) as e:
            st.error(f"Failed to parse research data from API response: {e}")
            return None
    return None

def get_ai_summary(text_to_summarize, api_key):
    """
    Calls the Gemini API to get a text summary.
    """
    st.info("🤖 Calling Generative AI for summarization...")
    
    prompt = f"Summarize the following text in one or two sentences, capturing the key takeaway: '{text_to_summarize}'"
    
    payload = {"contents": [{"role": "user", "parts": [{"text": prompt}]}]}
    result = call_gemini_api(payload, api_key)
    
    if result and 'candidates' in result:
        try:
            return result['candidates'][0]['content']['parts'][0]['text']
        except (KeyError, IndexError) as e:
            st.error(f"Failed to parse summary from API response: {e}")
            return "Summary could not be generated."
    return "Summary could not be generated."

# --- POWERPOINT GENERATION ---
def create_gtm_presentation(data):
    """
    Creates the final PowerPoint presentation based on user input.
    """
    prs = Presentation()
    
    # Use a simple title slide layout
    title_slide_layout = prs.slide_layouts[0]
    
    # --- Title Slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"GTM Strategy: {data['project_name']}"
    subtitle.text = f"Prepared for {data['region']}"

    # Use a title and content layout for subsequent slides
    content_slide_layout = prs.slide_layouts[1]
    
    # --- Objectives Slide ---
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Regional Objectives"
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Global Objectives Alignment:\n"
    p = tf.add_paragraph()
    p.text = data.get('global_objectives', 'Not found.')
    p.level = 1
    tf.add_paragraph() # Spacer
    tf.add_paragraph().text = "Regional Alignments & Additions:"
    p = tf.add_paragraph()
    p.text = data.get('regional_objectives', 'Not provided.')
    p.level = 1
    
    # --- Market Insights Slide ---
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = f"Market Insights: {data['region']}"
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Automated Deep Research Insights:"
    api_insights_data = data.get('api_insights', {})
    if api_insights_data:
        for key, value in api_insights_data.items():
            p = tf.add_paragraph()
            p.text = f"- {key.replace('_', ' ').title()}: {value}"
            p.level = 1
    else:
        p = tf.add_paragraph()
        p.text = "No research data generated."
        p.level = 1
        
    tf.add_paragraph() # Spacer
    p = tf.add_paragraph()
    p.text = "Custom Regional Insights:"
    p = tf.add_paragraph()
    p.text = data.get('custom_insights', 'Not provided.')
    p.level = 1
    
    # --- Timeline Slide ---
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Project Timeline"
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for item in data.get('timeline', []):
        p = tf.add_paragraph()
        p.text = item
    
    # --- Activation Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only Layout
    slide.shapes.title.text = "Regional Activation Plan"
    
    # Add three text boxes for the layout
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(3.0), Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height); tf = txBox.text_frame
    p = tf.add_paragraph(); p.text = "Insights Summary"; p.font.bold = True
    p = tf.add_paragraph(); p.text = data.get('activation_insights_summary', 'Not generated.')
    
    left = Inches(3.5)
    txBox = slide.shapes.add_textbox(left, top, width, height); tf = txBox.text_frame
    p = tf.add_paragraph(); p.text = "Activation Plan"; p.font.bold = True
    p = tf.add_paragraph(); p.text = data.get('activation_plan', 'Not provided.')

    left = Inches(6.5)
    txBox = slide.shapes.add_textbox(left, top, width, height); tf = txBox.text_frame
    p = tf.add_paragraph(); p.text = "Measurement & KPIs"; p.font.bold = True
    p = tf.add_paragraph(); p.text = data.get('measurement_plan', 'Not provided.')
    
    # --- Investment Slide ---
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Investment Summary"
    body_shape = slide.shapes.placeholders[1]
    
    try:
        investment_data = data.get('investment_data', [])
        if investment_data:
            df = pd.DataFrame(investment_data)
            rows, cols = df.shape
            table = slide.shapes.add_table(rows+1, cols, Inches(1.5), Inches(2.0), Inches(7), Inches(1.0)).table
            for col_idx, col_name in enumerate(df.columns): table.cell(0, col_idx).text = col_name
            for row_idx in range(rows):
                for col_idx in range(cols): table.cell(row_idx+1, col_idx).text = str(df.iloc[row_idx, col_idx])
        else:
            body_shape.text_frame.text = "No investment data provided."
    except Exception as e:
        body_shape.text_frame.text = f"Could not generate investment table. Please check format.\nError: {e}"

    # --- Overview Slide ---
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "AI-Generated Overview"
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = data.get('overview_summary', 'Not generated.')
    
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# --- SESSION STATE MANAGEMENT ---
if 'step' not in st.session_state:
    st.session_state.step = 0
    st.session_state.project_data = {}
    st.session_state.api_key = ""

def next_step(): st.session_state.step += 1
def prev_step(): st.session_state.step -= 1

# --- UI LAYOUT ---
st.sidebar.title("DreamAI Setups 🤖")
st.sidebar.markdown("---")
st.sidebar.header("Configuration")
api_key_input = st.sidebar.text_input("Enter Your Google AI API Key", type="password", key="api_key_input_widget")
if api_key_input:
    st.session_state.api_key = api_key_input
st.sidebar.markdown("---")


# --- Step 0: Welcome & Upload ---
if st.session_state.step == 0:
    st.title("Welcome to DreamAI Setups")
    st.markdown("Automate the creation of regional GTM slide decks. Start by uploading your global GTM presentation.")
    
    if not st.session_state.api_key:
        st.warning("Please enter your Google AI API Key in the sidebar to begin.")

    project_name = st.text_input("Enter a Project Name:", key="project_name_input")
    region = st.selectbox("Select Target Region:", ["Australia", "Japan", "Korea", "China"], key="region_input")
    uploaded_file = st.file_uploader("Upload Global GTM Deck (.pptx)", type="pptx", key="uploader", disabled=not st.session_state.api_key)

    if st.button("Start Analysis & Build", type="primary", disabled=not st.session_state.api_key):
        if uploaded_file and project_name and region:
            with st.spinner('Analyzing Presentation... This may take a moment.'):
                st.session_state.project_data['project_name'] = project_name
                st.session_state.project_data['region'] = region
                # For this example, we still mock the content parsing part.
                st.session_state.project_data['global_objectives'] = "Increase market share by 15% globally.\nLaunch Product X in 10 new markets.\nAchieve a 20% growth in online sales."
                st.session_state.project_data['global_timeline_points'] = ["Q1: Global Kick-off", "Q2: Product Finalization", "Q3: Marketing Campaign Launch", "Q4: Sales Push"]
                next_step()
                st.rerun()
        else:
            st.error("Please provide a project name, select a region, and upload a .pptx file.")

# --- Multi-Step Form ---
if st.session_state.step > 0:
    st.sidebar.header(f"Project: {st.session_state.project_data.get('project_name', 'New Project')}")
    st.sidebar.markdown(f"**Region:** {st.session_state.project_data.get('region', 'N/A')}")
    st.sidebar.markdown("---")
    
    progress_value = st.session_state.step / 6
    st.progress(progress_value)

    if st.sidebar.button("↩️ Start Over"):
        st.session_state.step = 0
        st.session_state.project_data = {}
        st.rerun()
    st.sidebar.markdown("---")

# --- Step 1: Objectives Alignment ---
if st.session_state.step == 1:
    st.header("Step 1: Objectives Alignment")
    st.subheader("Global Objectives Extracted")
    st.markdown("The following objectives were extracted from the global deck.")
    st.text_area("Global Objectives", value=st.session_state.project_data['global_objectives'], height=150, disabled=True)
    st.subheader("Regional Objectives Questionnaire")
    st.markdown("Confirm and adapt these objectives for your region. Add any region-specific goals.")
    regional_objectives = st.text_area("Enter your specific regional objectives:", height=200, key="regional_obj_input")
    if st.button("Save & Next: Insights →", type="primary", use_container_width=True):
        st.session_state.project_data['regional_objectives'] = regional_objectives
        next_step()
        st.rerun()

# --- Step 2: Regional Insight Generation ---
if st.session_state.step == 2:
    st.header("Step 2: Regional Insight Generation")
    with st.spinner("AI is conducting research..."):
        api_insights = get_deep_research(st.session_state.project_data['region'], st.session_state.api_key)
    if api_insights:
        st.subheader("Automated Deep Research")
        st.markdown("The AI has pulled the following market data.")
        st.json(api_insights)
    else:
        st.error("Could not fetch AI-powered insights. Please check your API key and try again.")
    st.subheader("Custom Regional Insights")
    st.markdown("Add your own qualitative findings, customer feedback, or specific market knowledge.")
    custom_insights = st.text_area("Add your custom insights here:", height=200, key="custom_insight_input")
    col1, col2, col3 = st.columns([1,1,1]);
    with col1: st.button("← Back: Objectives", on_click=prev_step, use_container_width=True)
    with col3:
        if st.button("Save & Next: Timeline →", type="primary", use_container_width=True):
            st.session_state.project_data['api_insights'] = api_insights
            st.session_state.project_data['custom_insights'] = custom_insights
            next_step()
            st.rerun()

# --- Step 3: Timeline Customization ---
if st.session_state.step == 3:
    st.header("Step 3: Timeline Customization")
    st.markdown("The global timeline has been imported. Add your region-specific activations and milestones.")
    timeline_items = st.session_state.project_data['global_timeline_points'].copy()
    if 'regional_events' not in st.session_state.project_data:
        st.session_state.project_data['regional_events'] = []
    st.subheader("Add Regional Timeline Events")
    for item in st.session_state.project_data['regional_events']: st.text(f"- {item}")
    new_event = st.text_input("Add a new timeline event (e.g., 'Q3: Local Influencer Campaign')")
    if st.button("Add Event"):
        if new_event:
            st.session_state.project_data['regional_events'].append(new_event)
            st.rerun()
    final_timeline = timeline_items + st.session_state.project_data['regional_events']
    st.subheader("Preview of Final Timeline")
    st.expander("Click to see the combined timeline").write(final_timeline)
    col1, col2, col3 = st.columns([1,1,1])
    with col1: st.button("← Back: Insights", on_click=prev_step, use_container_width=True)
    with col3:
        if st.button("Save & Next: Activations →", type="primary", use_container_width=True):
            st.session_state.project_data['timeline'] = final_timeline
            next_step()
            st.rerun()

# --- Step 4: Activation Planning ---
if st.session_state.step == 4:
    st.header("Step 4: Activation Planning")
    st.markdown("Plan your regional activation on a single slide. The AI will summarize the insights for you.")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.subheader("Insights Summary")
        with st.spinner("AI is summarizing..."):
            insights_summary = get_ai_summary(st.session_state.project_data.get('custom_insights', ''), st.session_state.api_key)
        st.info(insights_summary)
    with col2:
        st.subheader("Activation Plan")
        activation_plan = st.text_area("Detail your specific marketing activations.", height=300, key="activation_plan_input")
    with col3:
        st.subheader("Measurement (KPIs)")
        measurement_plan = st.text_area("Define how you will measure success.", height=300, key="measurement_plan_input")
    col_back, col_mid, col_next = st.columns([1,1,1])
    with col_back: st.button("← Back: Timeline", on_click=prev_step, use_container_width=True)
    with col_next:
        if st.button("Save & Next: Investment →", type="primary", use_container_width=True):
            st.session_state.project_data['activation_insights_summary'] = insights_summary
            st.session_state.project_data['activation_plan'] = activation_plan
            st.session_state.project_data['measurement_plan'] = measurement_plan
            next_step()
            st.rerun()

# --- Step 5: Investment Summary ---
if st.session_state.step == 5:
    st.header("Step 5: Investment Summary")
    st.markdown("Provide the regional investment details in CSV format.")
    st.info("Paste budget data below. The first line should be headers.\nExample:\nCategory,Q1 Budget,Q2 Budget\nMedia Spend,$50000,$75000")
    investment_data_str = st.text_area("Paste budget data here (CSV format):", height=200, key="investment_input")
    investment_data = []
    if investment_data_str:
        lines = investment_data_str.strip().split('\n')
        if len(lines) > 1:
            try:
                headers = [h.strip() for h in lines[0].split(',')]
                for line in lines[1:]:
                    values = [v.strip() for v in line.split(',')]
                    if len(values) == len(headers): investment_data.append(dict(zip(headers, values)))
            except Exception: st.warning("Could not parse data. Please check CSV format.")
    st.subheader("Investment Data Preview")
    if investment_data: st.dataframe(pd.DataFrame(investment_data))
    else: st.warning("No valid data entered yet.")
    col_back, col_mid, col_next = st.columns([1,1,1])
    with col_back: st.button("← Back: Activations", on_click=prev_step, use_container_width=True)
    with col_next:
        if st.button("Save & Next: Overview →", type="primary", use_container_width=True):
            st.session_state.project_data['investment_data'] = investment_data
            next_step()
            st.rerun()

# --- Step 6: AI-Powered Overview & Export ---
if st.session_state.step == 6:
    st.header("Step 6: Final Review & Export")
    st.balloons()
    st.subheader("AI-Generated Overview")
    st.markdown("The AI has created a summary based on all the information you provided.")
    full_text = f"Project: {st.session_state.project_data.get('project_name')}\nRegion: {st.session_state.project_data.get('region')}\nObjectives: {st.session_state.project_data.get('regional_objectives')}\nActivations: {st.session_state.project_data.get('activation_plan')}\nMeasurement: {st.session_state.project_data.get('measurement_plan')}"
    with st.spinner("AI is generating the final overview..."):
        overview_summary = get_ai_summary(full_text, st.session_state.api_key)
    st.text_area("Overview Summary", value=overview_summary, height=200, disabled=True)
    st.session_state.project_data['overview_summary'] = overview_summary
    st.subheader("Export Your Presentation")
    st.markdown("Your regional GTM deck is ready to download.")
    try:
        ppt_file = create_gtm_presentation(st.session_state.project_data)
        st.download_button(label="📥 Download PowerPoint (.pptx)", data=ppt_file, file_name=f"{st.session_state.project_data['project_name']}_Regional_GTM.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
    except Exception as e:
        st.error(f"An error occurred while generating the presentation: {e}")
    col_back, col_mid = st.columns([1,1])
    with col_back: st.button("← Back: Investment", on_click=prev_step, use_container_width=True)
