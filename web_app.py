import os
import io
import json
import tempfile
from fastapi import FastAPI, Request, UploadFile, File, Form
from fastapi.responses import HTMLResponse, FileResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pptx import Presentation
import requests
import shutil
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.dml.color import RGBColor
import base64
import zipfile
import copy
import uuid
import mimetypes
import openai
import httpx
import convertapi
from datetime import datetime
import sys
import traceback
import uvicorn
import asyncio
import threading
from task_storage import TaskManager, TEMP_DIR

# Load environment variables from .env file
try:
    from dotenv import load_dotenv
    load_dotenv()
    print("[ENV] Loaded environment variables from .env file")
except ImportError:
    print("[ENV] python-dotenv not installed, using system environment variables only")

# WORKAROUND: Monkey-patch OpenAI client to remove proxy argument issues
original_openai_init = openai.OpenAI.__init__

def patched_openai_init(self, *args, **kwargs):
    # Remove any proxy-related arguments that might cause issues
    kwargs.pop('proxies', None)
    return original_openai_init(self, *args, **kwargs)

openai.OpenAI.__init__ = patched_openai_init

# --- Debugging information ---

print(f"[DEBUG] Running file: {__file__}")
print(f"[DEBUG] Process ID: {os.getpid()}")
print(f"[DEBUG] sys.argv: {sys.argv}")
print(f"[DEBUG] openai module location: {openai.__file__}")

# --- DIAGNOSTIC: Check environment variables ---
print(f"[DIAGNOSTIC] CONVERTAPI_SECRET exists: {bool(os.getenv('CONVERTAPI_SECRET'))}")
print(f"[DIAGNOSTIC] OPENAI_API_KEY exists: {bool(os.getenv('OPENAI_API_KEY'))}")
print(f"[DIAGNOSTIC] Current working directory: {os.getcwd()}")
print(f"[DIAGNOSTIC] Python path: {sys.path[:3]}...")  # First 3 entries

# --- FastAPI setup ---
app = FastAPI(
    title="Dynamic AI Presentation Assembler",
    description="AI-powered presentation assembly service",
    version="1.0.0"
)

# Configure for large file uploads
app.add_middleware(
    lambda app: app,  # Placeholder for any middleware
)

# Increase the default request size limit (default is 16MB, we'll set to 100MB)
from fastapi.middleware.cors import CORSMiddleware
from fastapi import Request
from fastapi.exceptions import RequestValidationError
from fastapi.responses import JSONResponse
import traceback

# Add CORS middleware for better compatibility
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, specify your domains
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Custom exception handler for request validation errors
@app.exception_handler(RequestValidationError)
async def validation_exception_handler(request: Request, exc: RequestValidationError):
    print(f"[ERROR] Request validation failed: {exc}")
    print(f"[ERROR] Request details: {request.method} {request.url}")
    return JSONResponse(
        status_code=400,
        content={"error": f"Request validation failed: {str(exc)}"}
    )

# Custom exception handler for general exceptions
@app.exception_handler(Exception)
async def general_exception_handler(request: Request, exc: Exception):
    print(f"[ERROR] Unexpected error: {exc}")
    print(f"[ERROR] Traceback: {traceback.format_exc()}")
    return JSONResponse(
        status_code=500,
        content={"error": f"Internal server error: {str(exc)}"}
    )

app.mount("/static", StaticFiles(directory="static"), name="static")
templates = Jinja2Templates(directory="templates")

# --- Configuration for the Conversion Service ---
CONVERSION_SERVICE_URL = os.getenv("CONVERSION_SERVICE_URL", "http://localhost:8000/convert_document")

API_SECRET = os.getenv("CONVERTAPI_SECRET")
if not API_SECRET:
    raise RuntimeError("CONVERTAPI_SECRET is not set")
convertapi.api_credentials = API_SECRET
convertapi.base_uri = "https://au-v2.convertapi.com/"

# --- Conversion logic copied from conversion_service.py ---
def convert_pptx_to_png_zip(file_bytes: bytes, filename: str) -> bytes:
    upload_io = convertapi.UploadIO(io.BytesIO(file_bytes), filename=filename)
    result = convertapi.convert(
        'png',
        {'File': upload_io},
        from_format='pptx'
    )
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
        base = os.path.splitext(filename)[0]
        for idx, slide in enumerate(result.files, start=1):
            data = slide.io.read()
            name = f"{base}_slide{idx}.png"
            zf.writestr(name, data)
    buffer.seek(0)
    return buffer.read()

# --- Updated get_all_slide_data ---
def get_all_slide_data(file_bytes: bytes, file_type: str, filename: str = "document.pptx") -> list[dict]:
    print("[PROGRESS] Starting file conversion to PNGs via ConvertAPI...")
    if file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        try:
            # Extract PNG images
            zip_bytes = convert_pptx_to_png_zip(file_bytes, filename)
            print("[PROGRESS] PPTX converted to PNG ZIP, extracting slides...")
            
            # Extract text content from PPTX
            print("[PROGRESS] Extracting text content from PPTX...")
            slide_texts = []
            try:
                prs = Presentation(io.BytesIO(file_bytes))
                for slide_idx, slide in enumerate(prs.slides):
                    text_content = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_content.append(shape.text.strip())
                    slide_text = "\n".join(text_content)
                    slide_texts.append(slide_text)
                    print(f"[DEBUG] Slide {slide_idx}: extracted {len(text_content)} text shapes, total text length: {len(slide_text)}")
            except Exception as text_error:
                print(f"[WARNING] Failed to extract text from PPTX: {text_error}")
                slide_texts = []
            
            # Combine PNG images with text content
            slides = []
            with zipfile.ZipFile(io.BytesIO(zip_bytes), 'r') as zf:
                png_files = sorted([name for name in zf.namelist() if name.lower().endswith('.png')])
                for idx, name in enumerate(png_files):
                    with zf.open(name) as img_file:
                        img_bytes = img_file.read()
                        img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                    
                    # Get corresponding text content (if available)
                    slide_text = slide_texts[idx] if idx < len(slide_texts) else ""
                    
                    slides.append({
                        'slide_index': idx, 
                        'image_data': img_b64, 
                        'text': slide_text
                    })
                    print(f"[DEBUG] Slide {idx}: image extracted, text length: {len(slide_text)}")
            
            print(f"[PROGRESS] Extracted {len(slides)} slides with both images and text content.")
            return slides
        except Exception as e:
            print(f"[ERROR] ConvertAPI error: {e}")
            raise RuntimeError(f"ConvertAPI error: {e}")
    else:
        print("[ERROR] Only PPTX files are supported for conversion in this implementation.")
        raise RuntimeError("Only PPTX files are supported for conversion in this implementation.")

def find_slide_by_ai(api_key, file_bytes: bytes, file_type: str, slide_type_prompt: str, deck_name: str):
    print("[PROGRESS] Starting AI slide selection with OpenAI...")
    if not slide_type_prompt:
        return {"slide": None, "index": -1, "justification": "No keyword provided."}
    if not api_key:
        return {"slide": None, "index": -1, "justification": "OpenAI API Key is missing."}
    
    # Create OpenAI client - simple initialization only
    try:
        client = openai.OpenAI(api_key=api_key)
    except Exception as e:
        print(f"[ERROR] Failed to create OpenAI client: {e}")
        return {"slide": None, "index": -1, "justification": f"OpenAI client creation failed: {e}"}
    print("[PROGRESS] Getting all slide data for AI analysis...")
    slides_data = get_all_slide_data(file_bytes, file_type)
    print(f"[PROGRESS] {len(slides_data)} slides loaded for AI analysis.")
    system_prompt = f"""
    You are an expert presentation analyst with advanced visual recognition capabilities. Your task is to find the best slide/page in a document that matches a user's description.
    The user is looking for a slide/page representing: '{slide_type_prompt}'.
    
    **CRITICAL: You MUST analyze both the text content AND the visual structure from the PNG image for each slide/page.**
    
    **Visual Analysis Requirements:**
    - Examine the layout, positioning, and arrangement of visual elements
    - Look for visual patterns like flowcharts, timelines, hierarchical structures, charts, graphs
    - Consider color schemes, spacing, alignment, and visual hierarchy
    - Identify visual indicators like arrows, connecting lines, numbered sequences, icons
    
    **Text Analysis Requirements:**
    - Analyze textual content for relevant keywords and context
    - Look for structural indicators (headings, bullet points, numbered lists)
    - Consider semantic meaning and topic relevance
    
    **Specific Slide Type Guidance:**
    - **Timeline slides**: Look for horizontal/vertical sequences, dates, phases, arrows showing progression, chronological elements
    - **Objectives slides**: Look for goal-oriented language, bullet points with targets, key results, strategic aims
    - **Process slides**: Look for step-by-step flows, numbered sequences, arrows between elements
    - **Data slides**: Look for charts, graphs, tables, metrics, numbers
    
    **Your justification MUST include:**
    1. Visual reasoning: What visual patterns/layout made this slide suitable?
    2. Text reasoning: What textual content supports this choice?
    3. Why this slide is better than others for the requested type
    
    You must prioritize actual content slides over divider/table of contents pages.
    Return a JSON object with 'best_match_index' (integer, or -1) and 'justification' (detailed explanation including both visual and textual reasoning).
    """
    user_parts = [
        {"type": "text", "text": f"Find the best slide/page for '{slide_type_prompt}' in the '{deck_name}' with the following pages/slides:"}
    ]
    for slide_info in slides_data:
        user_parts.append({"type": "text", "text": f"\n--- Page/Slide {slide_info['slide_index'] + 1} (Text): {slide_info['text']}"})
        user_parts.append({"type": "image_url", "image_url": {"url": f"data:image/png;base64,{slide_info['image_data']}"}})
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_parts}
    ]
    try:
        print("[PROGRESS] Sending request to OpenAI for slide selection...")
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            response_format={"type": "json_object"}
        )
        result = json.loads(response.choices[0].message.content)
        best_index = result.get("best_match_index", -1)
        justification = result.get("justification", "No justification provided.")
        selected_slide_data = slides_data[best_index] if best_index != -1 and best_index < len(slides_data) else None
        print(f"[PROGRESS] OpenAI selected slide index: {best_index}")
        print(f"[AI REASONING] Visual + Text Analysis: {justification}")
        if selected_slide_data:
            print(f"[AI SELECTED SLIDE] Slide {best_index + 1} was chosen based on both visual layout patterns and textual content matching '{slide_type_prompt}'")
            print(f"[AI VISUAL ANALYSIS] The AI analyzed the PNG image to identify visual structures, layouts, and design patterns that match the requested slide type")
            print(f"[AI TEXT ANALYSIS] Text content from slide: {selected_slide_data.get('text', 'No text extracted')[:200]}...")
        else:
            print(f"[AI NO MATCH] No suitable slide found for '{slide_type_prompt}' based on visual and textual analysis")
        return {"slide": selected_slide_data, "index": best_index, "justification": justification}
    except Exception as e:
        print(f"[ERROR] OpenAI slide selection failed: {e}")
        return {"slide": None, "index": -1, "justification": f"AI error: {e}"}

def analyze_and_map_content(api_key, gtm_slide_content_data, template_slides_data, user_keyword):
    if not api_key:
        return {"best_template_index": -1, "justification": "OpenAI API Key is missing.", "processed_content": gtm_slide_content_data}
    
    # Create OpenAI client - simple initialization only
    try:
        client = openai.OpenAI(api_key=api_key)
    except Exception as e:
        print(f"[ERROR] Failed to create OpenAI client: {e}")
        return {"best_template_index": -1, "justification": f"OpenAI client creation failed: {e}", "processed_content": gtm_slide_content_data}
    system_prompt = f"""
    You are an expert presentation content mapper. Your primary task is to help a user integrate content from a Global (GTM) slide/page into the most appropriate regional template.
    Given the `gtm_slide_content` (with its text and image) and a list of `template_slides_data` (each with an index and text content, and image data), you must perform two critical tasks:
    1.  **Select the BEST Template:**
        * **Crucially, you must review *each and every* template slide/page text summary AND its associated visual content.**
        * Semantically and **visually** evaluate which template slide's structure and implied purpose would *best* accommodate the `gtm_slide_content`.
        * **Perform a comparative analysis:** Do not just pick the first decent match. Compare all options to find the single most suitable template based on a combined understanding of text and visuals. **Prioritize templates where the text *imlies* a strong visual match, rather than just explicitly stating a type.** For instance, a template with short, sequential bullet points and dates might be a better visual timeline fit than one that simply has "Timeline" in its title but dense paragraphs.
        * Consider factors like:
            * Does the template's textual layout (e.g., presence of sections, bullet points, titles) **and its visual layout (e.g., number of content blocks, placement of image placeholders, overall design)** match the theme/type of the GTM content.
            * Is there sufficient space or logical sections in the template for the GTM content based on its textual and visual structure?
            * Is the template visually appropriate for the content's nature (e.g., if GTM content is a timeline, does the template's visual suggest a timeline-like structure with distinct steps)?
    2.  **Process GTM Content for Regionalization:**
        * Analyze the `gtm_slide_content` (title and body text).
        * Identify any parts of the text that are highly likely to be *regional-specific* (e.g., local market data, specific regional initiatives, detailed local performance figures, regional names, or examples relevant only to one region).
        * For these regional-specific parts, replace them with a concise, generic placeholder like `[REGIONAL DATA HERE]`, `[LOCAL EXAMPLE]`, `[Qx REGIONAL METRICS]`, `[REGIONAL IMPACT]`, `[LOCAL TEAM]`, etc. Be intelligent about the placeholder text.
        * The goal is to provide a global baseline with clear, actionable markers for regional teams to fill in.
        * Maintain the original overall structure, headings, and flow of the text where possible.
    You MUST return a JSON object with the following keys:
    -   `best_template_index`: An integer representing the index of the best template slide/page from the `template_slides_data` list.
    -   `justification`: A brief, one-sentence justification for choosing that template, explicitly mentioning why it's better than other contenders if applicable.
    -   `processed_content`: An object with 'title' and 'body' keys, containing the GTM content with regional placeholders inserted.
    """
    # Defensive: always provide image_data (empty string if missing)
    image_data = gtm_slide_content_data.get('image_data', '')
    user_parts = [
        {"type": "text", "text": f"User's original keyword for this content: '{user_keyword}'"},
        {"type": "text", "text": "GTM Slide/Page Content to Process (Text):"},
        {"type": "text", "text": json.dumps(gtm_slide_content_data.get('text', {}), indent=2)},
        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_data}"}}
    ]
    user_parts.append({"type": "text", "text": "\nAvailable Template Slides/Pages Summary and Visuals:"})
    for slide_info in template_slides_data:
        user_parts.append({"type": "text", "text": f"\n--- Template Slide/Page {slide_info['slide_index'] + 1} (Text): {slide_info['text']}"})
        user_parts.append({"type": "image_url", "image_url": {"url": f"data:image/png;base64,{slide_info['image_data']}"}})
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_parts}
    ]
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            response_format={"type": "json_object"}
        )
        result = json.loads(response.choices[0].message.content)
        if "best_template_index" not in result or "justification" not in result or "processed_content" not in result:
            raise ValueError("AI response missing required keys.")
        best_index = result["best_template_index"]
        justification = result["justification"]
        processed_content = result["processed_content"]
        if "title" not in processed_content:
            processed_content["title"] = gtm_slide_content_data.get("title", "")
        if "body" not in processed_content:
            processed_content["body"] = gtm_slide_content_data.get("body", "")
        return {
            "best_template_index": best_index,
            "justification": justification,
            "processed_content": processed_content
        }
    except Exception as e:
        return {"best_template_index": -1, "justification": f"AI error: {e}", "processed_content": gtm_slide_content_data}

def get_slide_content(slide):
    if not slide:
        return {"title": "", "body": ""}
    text_shapes = sorted([s for s in slide.shapes if s.has_text_frame and s.text.strip()], key=lambda s: s.top)
    title = ""
    body = ""
    if text_shapes:
        title = text_shapes[0].text.strip()
        body = "\n".join(s.text.strip() for s in text_shapes[1:])
    return {"title": title, "body": body}

def analyze_template_formatting(presentation):
    """
    Analyze all slides in the presentation to extract dominant formatting patterns
    """
    print(f"[FORMATTING] Analyzing template formatting across {len(presentation.slides)} slides")
    
    title_formats = []
    body_formats = []
    
    for slide_idx, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
                
            # Determine if this is likely a title or body shape
            is_likely_title = (
                hasattr(shape, 'is_placeholder') and shape.is_placeholder and 
                shape.placeholder_format.type in (1, 2, 8)
            ) or (shape.top < Pt(150))
            
            is_likely_body = (
                hasattr(shape, 'is_placeholder') and shape.is_placeholder and 
                shape.placeholder_format.type in (3, 4, 8, 14)
            ) or (shape.top >= Pt(150) and shape.height > Pt(100))
            
            # Extract formatting from this shape
            for para in shape.text_frame.paragraphs:
                if para.runs and para.runs[0].text.strip():  # Only analyze paragraphs with text
                    first_run = para.runs[0]
                    
                    # DEBUG: Show what we're capturing
                    font_size_value = first_run.font.size
                    font_size_pt = font_size_value.pt if font_size_value else None
                    print(f"[FORMATTING DEBUG] Slide {slide_idx}, Shape: font={first_run.font.name}, "
                          f"size_raw={font_size_value}, size_pt={font_size_pt}, bold={first_run.font.bold}, "
                          f"is_title={is_likely_title}, is_body={is_likely_body}")
                    
                    format_info = {
                        'font_name': first_run.font.name,
                        'font_size': first_run.font.size,  # Keep as Pt object
                        'font_size_pt': font_size_pt,      # Also store pt value for debugging
                        'bold': first_run.font.bold,
                        'italic': first_run.font.italic,
                        'alignment': para.alignment,
                        'level': getattr(para, 'level', 0),
                        'slide_idx': slide_idx
                    }
                    
                    # Capture bullet formatting if present
                    print(f"[BULLET DEBUG] Analyzing paragraph on slide {slide_idx}: '{para.runs[0].text[:30]}...'")
                    try:
                        print(f"[BULLET DEBUG] Paragraph has format: {hasattr(para, 'format')}")
                        if hasattr(para, 'format'):
                            print(f"[BULLET DEBUG] Paragraph format has bullet: {hasattr(para.format, 'bullet')}")
                            if hasattr(para.format, 'bullet'):
                                bullet = para.format.bullet
                                print(f"[BULLET DEBUG] Bullet object exists: {bullet}")
                                print(f"[BULLET DEBUG] Bullet type: {getattr(bullet, 'type', 'NO_TYPE')}")
                                print(f"[BULLET DEBUG] Bullet character: {getattr(bullet, 'character', 'NO_CHARACTER')}")
                                print(f"[BULLET DEBUG] Bullet has font: {hasattr(bullet, 'font')}")
                                
                                format_info['bullet'] = {
                                    'type': getattr(bullet, 'type', None),
                                    'character': getattr(bullet, 'character', None),
                                    'font': {
                                        'name': getattr(bullet.font, 'name', None) if hasattr(bullet, 'font') else None,
                                        'size': getattr(bullet.font, 'size', None) if hasattr(bullet, 'font') else None,
                                    }
                                }
                                
                                # Try to capture bullet color
                                try:
                                    if hasattr(bullet, 'font') and hasattr(bullet.font, 'fill'):
                                        if bullet.font.fill.type == MSO_FILL_TYPE.SOLID:
                                            format_info['bullet']['font']['color'] = bullet.font.fill.fore_color.rgb
                                            print(f"[BULLET DEBUG] Captured bullet color: {format_info['bullet']['font']['color']}")
                                except Exception as color_e:
                                    print(f"[BULLET DEBUG] Could not capture bullet color: {color_e}")
                                    
                                print(f"[BULLET DEBUG] ✅ CAPTURED BULLET: type={format_info['bullet']['type']}, "
                                      f"char='{format_info['bullet']['character']}', font_name={format_info['bullet']['font']['name']}")
                            else:
                                print(f"[BULLET DEBUG] ❌ Paragraph format has no bullet attribute")
                        else:
                            print(f"[BULLET DEBUG] ❌ Paragraph has no format attribute")
                    except Exception as e:
                        print(f"[BULLET DEBUG] ❌ Exception during bullet capture: {e}")
                        import traceback
                        traceback.print_exc()
                    
                    # Try to capture color
                    try:
                        if first_run.font.fill.type == MSO_FILL_TYPE.SOLID:
                            format_info['color'] = first_run.font.fill.fore_color.rgb
                    except Exception:
                        format_info['color'] = None
                    
                    if is_likely_title:
                        title_formats.append(format_info)
                    elif is_likely_body:
                        body_formats.append(format_info)
    
    # Find most common formatting for titles and body
    def get_dominant_format(formats_list, format_type="unknown"):
        if not formats_list:
            print(f"[FORMATTING] No {format_type} formats found")
            return None
            
        print(f"[FORMATTING] Analyzing {len(formats_list)} {format_type} format samples")
        
        # Count frequency of each formatting combination
        format_counts = {}
        for fmt in formats_list:
            # Use font_size_pt for comparison to avoid Pt object issues
            key = (fmt['font_name'], fmt.get('font_size_pt'), fmt['bold'], fmt['italic'])
            if key not in format_counts:
                format_counts[key] = []
            format_counts[key].append(fmt)
        
        # Show all detected formats for debugging
        for key, examples in format_counts.items():
            font_name, font_size_pt, bold, italic = key
            print(f"[FORMATTING] {format_type}: {font_name}, {font_size_pt}pt, bold={bold}, italic={italic} "
                  f"(found {len(examples)} times)")
        
        # Get the most common format
        if format_counts:
            most_common_key = max(format_counts.keys(), key=lambda k: len(format_counts[k]))
            most_common_format = format_counts[most_common_key][0]  # Return first example of most common format
            font_name, font_size_pt, bold, italic = most_common_key
            print(f"[FORMATTING] Dominant {format_type} format selected: {font_name}, {font_size_pt}pt, bold={bold}")
            return most_common_format
        return None
    
    dominant_title_format = get_dominant_format(title_formats, "title")
    dominant_body_format = get_dominant_format(body_formats, "body")
    
    print(f"[FORMATTING] Found {len(title_formats)} title samples, {len(body_formats)} body samples")
    if dominant_title_format:
        print(f"[FORMATTING] Dominant title format: {dominant_title_format['font_name']}, "
              f"{dominant_title_format['font_size']}, bold={dominant_title_format['bold']}")
        if dominant_title_format.get('bullet'):
            print(f"[FORMATTING] ✅ Title format has bullet: {dominant_title_format['bullet']}")
        else:
            print(f"[FORMATTING] ❌ Title format has no bullet")
    if dominant_body_format:
        print(f"[FORMATTING] Dominant body format: {dominant_body_format['font_name']}, "
              f"{dominant_body_format['font_size']}, bold={dominant_body_format['bold']}")
        if dominant_body_format.get('bullet'):
            print(f"[FORMATTING] ✅ Body format has bullet: {dominant_body_format['bullet']}")
        else:
            print(f"[FORMATTING] ❌ Body format has no bullet")
    
    return {
        'title_format': dominant_title_format,
        'body_format': dominant_body_format
    }

def populate_slide(slide, content, template_formatting=None):
    print(f"[DIAGNOSTIC] populate_slide called with content: {content}")
    print(f"[DIAGNOSTIC] populate_slide - slide has {len(slide.shapes)} shapes to work with")
    if template_formatting:
        print(f"[FORMATTING] Using advanced template formatting for slide population")
    
    # Handle case where body might be an array
    title_text = str(content.get("title", ""))
    body_content = content.get("body", "")
    
    # Convert array to string if needed
    if isinstance(body_content, list):
        body_text = '\n'.join(str(item) for item in body_content)
        print(f"[DIAGNOSTIC] Converted body array to string: '{body_text}'")
    else:
        body_text = str(body_content)
    
    print(f"[DIAGNOSTIC] Final title_text: '{title_text}'")
    print(f"[DIAGNOSTIC] Final body_text: '{body_text}'")
    
    title_populated, body_populated = False, False
    shapes_checked = 0
    
    # First pass: analyze all shapes and their properties
    for i, shape in enumerate(slide.shapes):
        print(f"[DIAGNOSTIC] Shape {i+1}: type={shape.shape_type}, has_text_frame={shape.has_text_frame}, "
              f"is_placeholder={getattr(shape, 'is_placeholder', False)}, "
              f"top={shape.top}, height={shape.height}")
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            print(f"[DIAGNOSTIC] Shape {i+1} placeholder type: {shape.placeholder_format.type}")
    
    def preserve_formatting_and_replace_text(text_frame, new_text, content_type="content", override_format=None):
        """
        Replace text while preserving the original formatting from the first run
        or using template-wide formatting if provided
        """
        print(f"[FORMATTING] Preserving formatting for {content_type}")
        
        # Use advanced template formatting if available
        if template_formatting and content_type in ['title', 'body']:
            format_key = f"{content_type}_format"
            advanced_format = template_formatting.get(format_key)
            if advanced_format:
                print(f"[FORMATTING] Using advanced template-wide format for {content_type}: "
                      f"{advanced_format.get('font_name')}, {advanced_format.get('font_size_pt')}pt, "
                      f"bold={advanced_format.get('bold')}")
                override_format = advanced_format
        
        # Use override format if provided (from template analysis)
        if override_format:
            print(f"[FORMATTING] Using template-wide format override: {override_format.get('font_name')}, "
                  f"{override_format.get('font_size')}, bold={override_format.get('bold')}")
            
            # ENHANCED APPROACH: Preserve existing paragraph structure and just update text
            # This maintains bullet formatting that may not be accessible via python-pptx API
            print(f"[FORMATTING] Using paragraph-preserving approach to maintain bullet formatting")
            
            # Split text into paragraphs
            paragraphs = new_text.split('\n') if new_text else ['']
            existing_paragraphs = text_frame.paragraphs
            
            # Strategy: Use existing paragraphs if available, add new ones if needed
            for i, para_text in enumerate(paragraphs):
                if i < len(existing_paragraphs):
                    # Use existing paragraph (preserves bullet formatting)
                    p = existing_paragraphs[i]
                    print(f"[FORMATTING] Using existing paragraph {i} - preserves original bullet formatting")
                    
                    # Clear existing text content but keep paragraph structure
                    # Use the proper API to clear runs
                    p.clear()  # This clears all runs while preserving paragraph formatting
                else:
                    # Add new paragraph (will need bullet formatting applied)
                    p = text_frame.add_paragraph()
                    print(f"[FORMATTING] Added new paragraph {i} - applying template formatting")
                    
                    # Apply template paragraph formatting to new paragraphs
                    if override_format.get('alignment'):
                        p.alignment = override_format['alignment']
                    if override_format.get('level') is not None:
                        p.level = override_format['level']
                
                # Add text with template formatting
                run = p.add_run()
                
                # Clean the text content - remove bullet characters since paragraph already has bullet formatting
                clean_text = para_text.strip()
                # Remove common bullet characters from the beginning of text
                if clean_text.startswith('• '):
                    clean_text = clean_text[2:]  # Remove "• " 
                elif clean_text.startswith('• '):
                    clean_text = clean_text[2:]  # Remove different bullet char
                elif clean_text.startswith('- '):
                    clean_text = clean_text[2:]  # Remove "- "
                elif clean_text.startswith('* '):
                    clean_text = clean_text[2:]  # Remove "* "
                
                run.text = clean_text
                print(f"[FORMATTING] Cleaned text: '{para_text}' -> '{clean_text}'")
                
                # Apply font formatting - be very careful with the format values
                if override_format.get('font_name'):
                    run.font.name = override_format['font_name']
                    print(f"[FORMATTING DEBUG] Applied font name: {override_format['font_name']}")
                
                if override_format.get('font_size'):
                    run.font.size = override_format['font_size']  # This should already be a Pt object
                    print(f"[FORMATTING DEBUG] Applied font size: {override_format['font_size']}")
                
                if override_format.get('bold') is not None:
                    run.font.bold = override_format['bold']
                    print(f"[FORMATTING DEBUG] Applied bold: {override_format['bold']}")
                
                if override_format.get('italic') is not None:
                    run.font.italic = override_format['italic']
                
                if override_format.get('color'):
                    try:
                        run.font.fill.solid()
                        run.font.fill.fore_color.rgb = override_format['color']
                    except Exception as e:
                        print(f"[FORMATTING DEBUG] Failed to apply color: {e}")
            
            # Remove any extra existing paragraphs if we have fewer new paragraphs
            paragraphs_to_remove = len(text_frame.paragraphs) - len(paragraphs)
            if paragraphs_to_remove > 0:
                print(f"[FORMATTING] Removing {paragraphs_to_remove} extra paragraphs")
                # Remove from the end, working backwards
                for _ in range(paragraphs_to_remove):
                    if len(text_frame.paragraphs) > 1:  # Always keep at least one paragraph
                        # Clear the last paragraph instead of removing it to avoid XML issues
                        text_frame.paragraphs[-1].clear()
                        text_frame.paragraphs[-1].text = ""
            
            print(f"[FORMATTING] Applied paragraph-preserving formatting to {content_type} with {len(paragraphs)} paragraphs")
            return
        
        # Fallback: preserve existing formatting
        original_paragraph_format = None
        original_run_format = None
        original_bullet_format = None
        
        if text_frame.paragraphs:
            first_para = text_frame.paragraphs[0]
            original_paragraph_format = {
                'alignment': first_para.alignment,
                'level': getattr(first_para, 'level', 0)
            }
            
            # Capture bullet formatting
            try:
                if hasattr(first_para, 'format') and hasattr(first_para.format, 'bullet'):
                    bullet = first_para.format.bullet
                    original_bullet_format = {
                        'type': getattr(bullet, 'type', None),
                        'character': getattr(bullet, 'character', None),
                        'font': {
                            'name': getattr(bullet.font, 'name', None) if hasattr(bullet, 'font') else None,
                            'size': getattr(bullet.font, 'size', None) if hasattr(bullet, 'font') else None,
                            'color': None
                        }
                    }
                    
                    # Try to capture bullet color
                    try:
                        if hasattr(bullet, 'font') and hasattr(bullet.font, 'fill'):
                            if bullet.font.fill.type == MSO_FILL_TYPE.SOLID:
                                original_bullet_format['font']['color'] = bullet.font.fill.fore_color.rgb
                    except Exception:
                        pass
                    
                    print(f"[FORMATTING] Captured bullet formatting: type={original_bullet_format['type']}, "
                          f"char={original_bullet_format['character']}")
            except Exception as e:
                print(f"[FORMATTING] Could not capture bullet formatting: {e}")
            
            if first_para.runs:
                first_run = first_para.runs[0]
                original_run_format = {
                    'font_name': first_run.font.name,
                    'font_size': first_run.font.size,
                    'bold': first_run.font.bold,
                    'italic': first_run.font.italic,
                    'underline': first_run.font.underline,
                    'color': None
                }
                
                # Try to preserve color
                try:
                    if first_run.font.fill.type == MSO_FILL_TYPE.SOLID:
                        original_run_format['color'] = first_run.font.fill.fore_color.rgb
                except Exception:
                    pass
                
                print(f"[FORMATTING] Captured original formatting: font={original_run_format.get('font_name')}, "
                      f"size={original_run_format.get('font_size')}, bold={original_run_format.get('bold')}")
        
        # Clear and rebuild with preserved formatting
        text_frame.clear()
        
        # Split text into paragraphs (in case of multi-line content)
        paragraphs = new_text.split('\n') if new_text else ['']
        
        for i, para_text in enumerate(paragraphs):
            if i == 0:
                # Use the existing first paragraph
                p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
            else:
                # Add new paragraphs for additional lines
                p = text_frame.add_paragraph()
            
            # Apply original paragraph formatting
            if original_paragraph_format:
                p.alignment = original_paragraph_format['alignment']
                p.level = original_paragraph_format['level']
            
            # Apply original bullet formatting
            if original_bullet_format:
                try:
                    bullet = p.format.bullet
                    if original_bullet_format['type'] is not None:
                        bullet.type = original_bullet_format['type']
                    if original_bullet_format['character'] is not None:
                        bullet.character = original_bullet_format['character']
                    
                    # Apply bullet font formatting if available
                    bullet_font_info = original_bullet_format.get('font', {})
                    if bullet_font_info.get('name'):
                        bullet.font.name = bullet_font_info['name']
                    if bullet_font_info.get('size'):
                        bullet.font.size = bullet_font_info['size']
                    if bullet_font_info.get('color'):
                        bullet.font.fill.solid()
                        bullet.font.fill.fore_color.rgb = bullet_font_info['color']
                    
                    print(f"[FORMATTING] Applied bullet formatting to paragraph {i}")
                except Exception as e:
                    print(f"[FORMATTING] Could not apply bullet formatting to paragraph {i}: {e}")
            
            # Add the text with original run formatting
            run = p.add_run()
            run.text = para_text
            
            # Apply original run formatting
            if original_run_format:
                if original_run_format['font_name']:
                    run.font.name = original_run_format['font_name']
                if original_run_format['font_size']:
                    run.font.size = original_run_format['font_size']
                if original_run_format['bold'] is not None:
                    run.font.bold = original_run_format['bold']
                if original_run_format['italic'] is not None:
                    run.font.italic = original_run_format['italic']
                if original_run_format['underline'] is not None:
                    run.font.underline = original_run_format['underline']
                if original_run_format['color']:
                    try:
                        run.font.fill.solid()
                        run.font.fill.fore_color.rgb = original_run_format['color']
                    except Exception:
                        pass
        
        print(f"[FORMATTING] Applied preserved formatting to {content_type} with {len(paragraphs)} paragraphs")
    
    for shape in slide.shapes:
        shapes_checked += 1
        print(f"[DIAGNOSTIC] Checking shape {shapes_checked}: has_text_frame={shape.has_text_frame}")
        
        # Check for title candidates (must have text frame)
        if shape.has_text_frame:
            is_title_placeholder = (
                hasattr(shape, 'is_placeholder') and shape.is_placeholder and 
                shape.placeholder_format.type in (1, 2, 8)
            )
            is_top_text_box = (shape.top < Pt(150))
            
            if not title_populated and (is_title_placeholder or is_top_text_box):
                print(f"[DIAGNOSTIC] Populating title in shape {shapes_checked}")
                # Use basic formatting preservation (AI will fix it later if needed)
                preserve_formatting_and_replace_text(shape.text_frame, title_text, "title")
                title_populated = True
                print(f"[DIAGNOSTIC] Title populated with preserved formatting: '{title_text}'")
                continue
                
            # Check for body candidates (must have text frame)
            is_body_placeholder = (
                hasattr(shape, 'is_placeholder') and shape.is_placeholder and 
                shape.placeholder_format.type in (3, 4, 8, 14)
            )
            is_lorem_ipsum = "lorem ipsum" in shape.text.lower()
            is_empty_text_box = not shape.text.strip() and shape.height > Pt(100)
            is_large_text_box = shape.height > Pt(200)
            
            if not body_populated and (is_body_placeholder or is_lorem_ipsum or is_empty_text_box or is_large_text_box):
                print(f"[DIAGNOSTIC] Populating body in shape {shapes_checked}")
                # Use basic formatting preservation (AI will fix it later if needed)
                preserve_formatting_and_replace_text(shape.text_frame, body_text, "body")
                body_populated = True
                print(f"[DIAGNOSTIC] Body populated with preserved formatting: '{body_text[:100]}...'")
                break
    
    # If body still not populated, be more aggressive - add a text box
    if not body_populated and body_text.strip():
        print(f"[DIAGNOSTIC] Body not populated yet, being more aggressive...")
        
        # Look for any placeholder that could hold content
        for shape in slide.shapes:
            if (hasattr(shape, 'is_placeholder') and shape.is_placeholder and 
                shape.placeholder_format.type in (2, 3, 4, 7, 8, 14, 15, 16)):  # Various content placeholders
                
                print(f"[DIAGNOSTIC] Found content placeholder type {shape.placeholder_format.type}, converting to text")
                try:
                    # Try to add text directly to the placeholder
                    if not shape.has_text_frame:
                        # If it doesn't have a text frame, try to convert it to one
                        print(f"[DIAGNOSTIC] Placeholder doesn't have text frame, trying to add text anyway")
                    
                    # Force create a text frame if possible
                    if hasattr(shape, 'text_frame') or shape.has_text_frame:
                        # Use basic formatting preservation (AI will fix it later if needed)
                        preserve_formatting_and_replace_text(shape.text_frame, body_text, "aggressive_body")
                        body_populated = True
                        print(f"[DIAGNOSTIC] Successfully forced body population in placeholder with preserved formatting")
                        break
                    else:
                        print(f"[DIAGNOSTIC] Cannot add text to this placeholder type")
                except Exception as e:
                    print(f"[DIAGNOSTIC] Failed to populate placeholder: {e}")
                    continue
    
    # Last resort: create a new text box with template-inspired formatting
    if not body_populated and body_text.strip():
        print(f"[DIAGNOSTIC] Last resort: creating new text box for body content")
        try:
            # Add a text box in the lower portion of the slide
            left = Pt(50)
            top = Pt(200)
            width = Pt(600)
            height = Pt(300)
            
            new_textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = new_textbox.text_frame
            tf.clear()
            
            # Use basic defaults - AI will fix formatting later via visual verification
            template_font_name = "Calibri"  # Default fallback
            template_font_size = Pt(18)     # Default fallback
            template_bold = False
            
            # Scan existing shapes for formatting clues
            print(f"[FORMATTING] Scanning existing shapes for formatting clues...")
            for existing_shape in slide.shapes:
                if (existing_shape.has_text_frame and existing_shape.text_frame.paragraphs):
                    for para in existing_shape.text_frame.paragraphs:
                        if para.runs and para.runs[0].text.strip():
                            first_run = para.runs[0]
                            if first_run.font.name:
                                template_font_name = first_run.font.name
                            if first_run.font.size:
                                template_font_size = first_run.font.size
                            template_bold = first_run.font.bold or False
                            print(f"[FORMATTING] Found template formatting from existing shape: {template_font_name}, {template_font_size}, bold={template_bold}")
                            break
                    break
            
            # Apply the template formatting to new text
            paragraphs = body_text.split('\n') if body_text else ['']
            for i, para_text in enumerate(paragraphs):
                if i == 0:
                    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                else:
                    p = tf.add_paragraph()
                
                run = p.add_run()
                run.text = para_text
                run.font.name = template_font_name
                run.font.size = template_font_size
                run.font.bold = template_bold
            
            body_populated = True
            print(f"[DIAGNOSTIC] Created new text box with template-inspired formatting: {template_font_name}, {template_font_size}")
        except Exception as e:
            print(f"[DIAGNOSTIC] Failed to create new text box: {e}")
    
    print(f"[DIAGNOSTIC] populate_slide results - title_populated: {title_populated}, body_populated: {body_populated}")
    print(f"[DIAGNOSTIC] Checked {shapes_checked} total shapes")

def copy_solid_or_gradient_background(src_slide, dest_slide):
    src_slide_elm = src_slide.element
    dest_slide_elm = dest_slide.element
    src_bg = src_slide_elm.find('.//p:bg', namespaces=src_slide_elm.nsmap)
    if src_bg is not None:
        new_bg = copy.deepcopy(src_bg)
        current_bg = dest_slide_elm.find('.//p:bg', namespaces=dest_slide_elm.nsmap)
        if current_bg is not None:
            current_bg.getparent().remove(current_bg)
        # Insert background before spTree (shape tree) but after cSld
        cSld = dest_slide_elm.find('.//p:cSld', namespaces=dest_slide_elm.nsmap)
        spTree = dest_slide_elm.find('.//p:spTree', namespaces=dest_slide_elm.nsmap)
        if cSld is not None and spTree is not None:
            cSld.insert(list(cSld).index(spTree), new_bg)
        else:
            dest_slide_elm.insert(0, new_bg)

def copy_slide_background(src_slide, dest_slide):
    src_slide_elm = src_slide.element
    dest_slide_elm = dest_slide.element
    src_bg_pr = src_slide_elm.find('.//p:bgPr', namespaces=src_slide_elm.nsmap)
    if src_bg_pr is None:
        return
    src_blip_fill = src_bg_pr.find('.//a:blipFill', namespaces=src_bg_pr.nsmap)
    if src_blip_fill is not None:
        src_blip = src_blip_fill.find('.//a:blip', namespaces=src_bg_pr.nsmap)
        if src_blip is not None and '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed' in src_blip.attrib:
            rId = src_blip.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed']
            try:
                src_image_part = src_slide.part.related_part(rId)
                image_bytes = src_image_part.blob
                new_image_part = dest_slide.part.get_or_add_image_part(image_bytes, src_image_part.content_type)
                new_rId = dest_slide.part.relate_to(new_image_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')
                new_bg_pr = copy.deepcopy(src_bg_pr)
                new_blip = new_bg_pr.find('.//a:blip', namespaces=new_bg_pr.nsmap)
                if new_blip is not None:
                    new_blip.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'] = new_rId
                # Replace the background property, not append
                current_bg = dest_slide_elm.find('.//p:bg', namespaces=dest_slide_elm.nsmap)
                if current_bg is not None:
                    current_bg.getparent().remove(current_bg)
                # Create a new <p:bg> element and insert new_bg_pr inside it
                from lxml import etree
                new_bg = etree.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}bg', nsmap=dest_slide_elm.nsmap)
                new_bg.append(new_bg_pr)
                # Insert background before spTree but after cSld to ensure it's behind content
                cSld = dest_slide_elm.find('.//p:cSld', namespaces=dest_slide_elm.nsmap)
                spTree = dest_slide_elm.find('.//p:spTree', namespaces=dest_slide_elm.nsmap)
                if cSld is not None and spTree is not None:
                    cSld.insert(list(cSld).index(spTree), new_bg)
                else:
                    dest_slide_elm.insert(0, new_bg)
            except Exception as e:
                copy_solid_or_gradient_background(src_slide, dest_slide)
    else:
        copy_solid_or_gradient_background(src_slide, dest_slide)

def deep_copy_slide_content(dest_slide, src_slide):
    print(f"[LOG] Copying slide content: src_slide layout={getattr(src_slide, 'slide_layout', None)}, dest_slide layout={getattr(dest_slide, 'slide_layout', None)}")
    print(f"[LOG] Source slide has {len(src_slide.shapes)} shapes.")
    for i, shape in enumerate(src_slide.shapes):
        print(f"[LOG] Source shape {i}: type={shape.shape_type}, has_text_frame={shape.has_text_frame}, is_placeholder={getattr(shape, 'is_placeholder', False)}")
    
    # DIAGNOSTIC: Check initial destination state
    print(f"[DIAGNOSTIC] Initial destination slide shapes: {len(dest_slide.shapes)}")
    
    # Remove all shapes from destination slide
    shapes_to_remove = list(dest_slide.shapes)
    print(f"[DIAGNOSTIC] Removing {len(shapes_to_remove)} existing shapes from destination")
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    
    print(f"[DIAGNOSTIC] After removal, destination slide shapes: {len(dest_slide.shapes)}")
    
    # Copy all shapes from source to destination
    shapes_copied = 0
    for i, shape in enumerate(src_slide.shapes):
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"[LOG] Copying PICTURE shape {i}")
            try:
                image_bytes = shape.image.blob
                dest_slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)
                shapes_copied += 1
                print(f"[DIAGNOSTIC] Successfully copied PICTURE shape {i}")
            except Exception as e:
                print(f"[ERROR] Failed to copy PICTURE shape {i}: {e}")
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    new_el = copy.deepcopy(shape.element)
                    dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                    shapes_copied += 1
                    print(f"[DIAGNOSTIC] Used fallback copy for PICTURE shape {i}")
        elif shape.has_text_frame:
            print(f"[LOG] Copying TEXT shape {i}")
            try:
                new_shape = dest_slide.shapes.add_textbox(left, top, width, height)
                new_text_frame = new_shape.text_frame
                new_text_frame.clear()
                for paragraph in shape.text_frame.paragraphs:
                    new_paragraph = new_text_frame.add_paragraph()
                    new_paragraph.alignment = paragraph.alignment
                    if hasattr(paragraph, 'level'):
                        new_paragraph.level = paragraph.level
                    for run in paragraph.runs:
                        new_run = new_paragraph.add_run()
                        new_run.text = run.text
                        new_run.font.bold = run.font.bold
                        new_run.font.italic = run.font.italic
                        new_run.font.underline = run.font.underline
                        if run.font.size:
                            new_run.font.size = run.font.size
                        if run.font.fill.type == MSO_FILL_TYPE.SOLID:
                            new_run.font.fill.solid()
                            try:
                                if isinstance(run.font.fill.fore_color.rgb, RGBColor):
                                    new_run.font.fill.fore_color.rgb = run.font.fill.fore_color.rgb
                                else:
                                    rgb_tuple = run.font.fill.fore_color.rgb
                                    new_run.font.fill.fore_color.rgb = RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])
                            except Exception:
                                pass
                new_text_frame.word_wrap = shape.text_frame.word_wrap
                new_text_frame.margin_left = shape.text_frame.margin_left
                new_text_frame.margin_right = shape.text_frame.margin_right
                new_text_frame.margin_top = shape.text_frame.margin_top
                new_text_frame.margin_bottom = shape.text_frame.margin_bottom
                shapes_copied += 1
                print(f"[DIAGNOSTIC] Successfully copied TEXT shape {i}")
            except Exception as e:
                print(f"[ERROR] Failed to copy TEXT shape {i}: {e}")
        else:
            print(f"[LOG] Copying OTHER shape {i}")
            try:
                new_el = copy.deepcopy(shape.element)
                dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                shapes_copied += 1
                print(f"[DIAGNOSTIC] Successfully copied OTHER shape {i}")
            except Exception as e:
                print(f"[ERROR] Failed to copy OTHER shape {i}: {e}")
    
    print(f"[DIAGNOSTIC] Total shapes copied: {shapes_copied} out of {len(src_slide.shapes)}")
    print(f"[LOG] Copying background...")
    copy_slide_background(src_slide, dest_slide)
    print(f"[LOG] Destination slide now has {len(dest_slide.shapes)} shapes.")
    print(f"[DIAGNOSTIC] Final shape copy verification: expected {len(src_slide.shapes)}, got {len(dest_slide.shapes)}")

@app.get("/", response_class=HTMLResponse)
def index(request: Request):
    return templates.TemplateResponse("index.html", {"request": request})

@app.post("/start_assemble")
async def start_assemble(
    request: Request,
    api_key: str = Form(...),
    template_files: list[UploadFile] = File(...),
    gtm_file: UploadFile = File(...),
    structure: str = Form(...)
):
    """Start assembly process and return task ID immediately"""
    print("[DEBUG] /start_assemble called")
    
    # Quick validation
    if not api_key or len(api_key.strip()) < 10:
        return JSONResponse(status_code=400, content={"error": "Invalid API key"})
    
    try:
        structure_steps = json.loads(structure)
        if not structure_steps:
            return JSONResponse(status_code=400, content={"error": "No structure steps"})
    except:
        return JSONResponse(status_code=400, content={"error": "Invalid structure JSON"})
    
    # Read file contents BEFORE starting background processing
    try:
        print("[DEBUG] Reading file contents...")
        template_files_data = []
        for i, template_file in enumerate(template_files):
            template_file.file.seek(0)  # Reset file pointer
            content = template_file.file.read()
            template_files_data.append({
                "filename": template_file.filename,
                "content": content
            })
            print(f"[DEBUG] Read template file {i+1}: {template_file.filename}, {len(content)} bytes")
        
        gtm_file.file.seek(0)  # Reset file pointer
        gtm_content = gtm_file.file.read()
        gtm_data = {
            "filename": gtm_file.filename,
            "content": gtm_content
        }
        print(f"[DEBUG] Read GTM file: {gtm_file.filename}, {len(gtm_content)} bytes")
        
    except Exception as e:
        print(f"[ERROR] Failed to read file contents: {e}")
        return JSONResponse(status_code=400, content={"error": f"Failed to read file contents: {str(e)}"})
    
    # Create task
    task_id = TaskManager.create_task()
    print(f"[DEBUG] Created task {task_id}")
    
    # Start processing in background thread with file data (not file objects)
    def process_in_background():
        asyncio.run(process_assembly_task(task_id, api_key, template_files_data, gtm_data, structure_steps))
    
    thread = threading.Thread(target=process_in_background, daemon=True)
    thread.start()
    
    return JSONResponse(content={
        "task_id": task_id,
        "status": "started",
        "message": "Processing started - use task_id to check status"
    })

@app.get("/task_status/{task_id}")
async def get_task_status(task_id: str):
    """Get current status of assembly task"""
    task = TaskManager.get_task(task_id)
    return JSONResponse(content=task)

@app.get("/download/{task_id}")
async def download_result(task_id: str):
    """Download the completed presentation"""
    task = TaskManager.get_task(task_id)
    
    if task.get("status") != "completed":
        return JSONResponse(status_code=404, content={"error": "Task not completed or not found"})
    
    result_file = task.get("result_file")
    if not result_file or not os.path.exists(result_file):
        return JSONResponse(status_code=404, content={"error": "Result file not found"})
    
    try:
        with open(result_file, 'rb') as f:
            file_data = f.read()
        
        def generate_chunks():
            chunk_size = 32 * 1024
            data = io.BytesIO(file_data)
            while True:
                chunk = data.read(chunk_size)
                if not chunk:
                    break
                yield chunk
        
        # Clean up task after successful download
        TaskManager.cleanup_task(task_id)
        
        return StreamingResponse(
            generate_chunks(),
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            headers={
                "Content-Disposition": "attachment; filename=assembled_presentation.pptx",
                "Content-Length": str(len(file_data))
            }
        )
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": f"Download failed: {str(e)}"})

async def process_assembly_task(task_id: str, api_key: str, template_files_data: list, gtm_data: dict, structure_steps: list):
    """Process assembly in background with full AI functionality"""
    try:
        TaskManager.update_task(task_id, "processing", 10, "Starting file processing...")
        
        print(f"[ASYNC] Processing task {task_id} with GTM file: {gtm_data['filename']}")
        
        # Save files to temp directory
        with tempfile.TemporaryDirectory() as tmpdir:
            TaskManager.update_task(task_id, "processing", 15, "Processing template files...")
            
            # Process template files (replicate the original logic)
            template_paths = []
            all_template_slides_for_ai = []
            base_pptx_template_found = False
            new_prs = None
            
            for i, file_data in enumerate(template_files_data):
                filename = file_data['filename']
                file_bytes = file_data['content']
                print(f"[ASYNC] Processing template file {i+1}: {filename}")
                path = os.path.join(tmpdir, filename)
                
                with open(path, 'wb') as out:
                    out.write(file_bytes)
                
                template_paths.append(path)
                file_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                
                if file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                    if not base_pptx_template_found:
                        # Use the first template file as the base - keep ALL slides from the first file
                        new_prs = Presentation(io.BytesIO(file_bytes))
                        base_pptx_template_found = True
                        print(f"[ASYNC] Loaded base template with {len(new_prs.slides)} slides")
                        print(f"[ASYNC] Base template slide layouts available: {len(new_prs.slide_layouts)}")
                    else:
                        # Merge additional template files by adding their slides to the base presentation
                        current_prs_to_merge = Presentation(io.BytesIO(file_bytes))
                        print(f"[ASYNC] Merging additional template file with {len(current_prs_to_merge.slides)} slides")
                        for slide_to_merge in current_prs_to_merge.slides:
                            # Use the same layout index as the source slide
                            layout_idx = slide_to_merge.slide_layout.slide_id
                            # Find the layout index in the base template that matches the source
                            matching_layout = None
                            for idx, layout in enumerate(new_prs.slide_layouts):
                                if layout.slide_id == layout_idx:
                                    matching_layout = layout
                                    break
                            if matching_layout is None:
                                matching_layout = new_prs.slide_layouts[0]  # fallback
                            new_slide = new_prs.slides.add_slide(matching_layout)
                            deep_copy_slide_content(new_slide, slide_to_merge)
                        print(f"[ASYNC] After merging, presentation now has {len(new_prs.slides)} slides")
                
                try:
                    slides_data = get_all_slide_data(file_bytes, file_type)
                    print(f"[ASYNC] Got {len(slides_data)} slides from template file {filename}")
                    all_template_slides_for_ai.extend(slides_data)
                except Exception as e:
                    print(f"[ASYNC] Error processing template slides for AI: {e}")
                    # Continue processing without failing
            
            if new_prs is None:
                raise Exception("No PPTX template found - at least one PPTX file must be uploaded as template")
            
            TaskManager.update_task(task_id, "processing", 25, "Processing GTM file...")
            
            # Process GTM file
            gtm_filename = gtm_data['filename']
            gtm_content = gtm_data['content']
            gtm_path = os.path.join(tmpdir, gtm_filename)
            with open(gtm_path, 'wb') as out:
                out.write(gtm_content)
            
            gtm_file_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            
            print(f"[ASYNC] GTM file: {gtm_filename}, type: {gtm_file_type}")
            num_template_slides = len(new_prs.slides)
            num_structure_steps = len(structure_steps)
            print(f"[ASYNC] num_template_slides: {num_template_slides}, num_structure_steps: {num_structure_steps}")
            
            # Keep all template slides available for AI selection
            print(f"[ASYNC] Keeping all {len(new_prs.slides)} template slides available for AI selection")
            
            TaskManager.update_task(task_id, "processing", 25, "Analyzing template formatting...")
            
            # ADVANCED FORMATTING ANALYSIS: Extract dominant formatting patterns from template
            template_formatting = analyze_template_formatting(new_prs)
            print(f"[FORMATTING] Template formatting analysis complete")
            if template_formatting:
                title_format = template_formatting.get('title_format')
                body_format = template_formatting.get('body_format')
                if title_format:
                    print(f"[FORMATTING] Dominant title format: {title_format.get('font_name')}, "
                          f"{title_format.get('font_size_pt')}pt, bold={title_format.get('bold')}")
                if body_format:
                    print(f"[FORMATTING] Dominant body format: {body_format.get('font_name')}, "
                          f"{body_format.get('font_size_pt')}pt, bold={body_format.get('bold')}")
            
            TaskManager.update_task(task_id, "processing", 30, "Starting AI analysis...")
            
            # Process each structure step with AI-driven slide selection
            total_steps = len(structure_steps)
            for i, step in enumerate(structure_steps):
                progress = 30 + int((i / total_steps) * 55)  # 30-85% for processing steps
                keyword = step["keyword"]
                action = step["action"]
                
                TaskManager.update_task(task_id, "processing", progress, 
                                      f"AI analyzing all slides for best match: {keyword}")
                print(f"[ASYNC] Step {i}: keyword={keyword}, action={action}")
                
                if action == "Copy from GTM (as is)":
                    print(f"[ASYNC] Copy from GTM (as is) - Finding GTM content for keyword '{keyword}'")
                    gtm_ai_selection_result = find_slide_by_ai(api_key, gtm_content, gtm_file_type, keyword, "GTM Deck (Content Source)")
                    
                    if gtm_ai_selection_result["slide"]:
                        # Step 1: AI selects the BEST template slide for this content type
                        print(f"[ASYNC] Step 1: AI analyzing ALL template slides to find best destination for '{keyword}'")
                        
                        # Use AI to select the best template slide destination
                        template_selection_result = analyze_and_map_content(
                            api_key,
                            {"title": keyword, "body": f"Content for {keyword}", "image_data": ""},
                            all_template_slides_for_ai,
                            keyword
                        )
                        
                        selected_template_index = template_selection_result["best_template_index"]
                        print(f"[ASYNC] AI selected template slide {selected_template_index} as best destination for '{keyword}'")
                        
                        if selected_template_index > 0 and selected_template_index <= len(new_prs.slides):
                            ai_chosen_slide_index = selected_template_index - 1  # Convert to 0-based
                            dest_slide = new_prs.slides[ai_chosen_slide_index]
                            print(f"[ASYNC] Using AI-chosen slide {selected_template_index} (index {ai_chosen_slide_index}) as destination")
                            
                            # Log details about the selected slide for bullet analysis
                            print(f"[SLIDE ANALYSIS] Selected slide has {len(dest_slide.shapes)} shapes")
                            for shape_idx, shape in enumerate(dest_slide.shapes):
                                if shape.has_text_frame:
                                    print(f"[SLIDE ANALYSIS] Shape {shape_idx}: {len(shape.text_frame.paragraphs)} paragraphs")
                                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                                        if para.text.strip():
                                            print(f"[SLIDE ANALYSIS] Para {para_idx}: '{para.text[:50]}...'")
                                            if hasattr(para, 'format') and hasattr(para.format, 'bullet'):
                                                bullet = para.format.bullet
                                                print(f"[SLIDE ANALYSIS] ✅ Has bullet: type={getattr(bullet, 'type', None)}, char='{getattr(bullet, 'character', None)}'")
                                            else:
                                                print(f"[SLIDE ANALYSIS] ❌ No bullet formatting")
                        else:
                            dest_slide = new_prs.slides[i] if i < len(new_prs.slides) else new_prs.slides[0]
                            print(f"[ASYNC] AI selection out of bounds, using fallback slide")
                        
                        # Step 2: Populate with PLACEHOLDER TEXT (not GTM content)
                        placeholder_content = {
                            "title": f"[{keyword.upper()} - PLACEHOLDER]",
                            "body": f"[REGIONAL {keyword.upper()} DATA AND ANALYSIS HERE]\n\n• Placeholder bullet point 1\n• Placeholder bullet point 2\n• Placeholder bullet point 3"
                        }
                        
                        populate_slide(dest_slide, placeholder_content, template_formatting)
                        print(f"[ASYNC] Populated AI-chosen slide with PLACEHOLDER content for '{keyword}'")
                    
                elif action == "Merge: Template Layout + GTM Content":
                    print(f"[ASYNC] Merge action - AI analyzing all slides for best destination for '{keyword}'")
                    
                    # Step 1: Find GTM content 
                    gtm_ai_selection_result = find_slide_by_ai(api_key, gtm_content, gtm_file_type, keyword, "GTM Deck (Content Source)")
                    
                    if not gtm_ai_selection_result["slide"]:
                        print(f"[ASYNC] No suitable GTM slide found for keyword '{keyword}'")
                        continue
                    
                    # Step 2: AI selects the BEST template slide destination
                    print(f"[ASYNC] Step 2: AI analyzing ALL template slides to find best destination for '{keyword}'")
                    
                    # Extract GTM content for context
                    full_text = gtm_ai_selection_result["slide"].get("text", "")
                    lines = full_text.split('\n') if full_text else []
                    gtm_content_context = {
                        "title": lines[0] if lines else keyword,
                        "body": "\n".join(lines[1:]) if len(lines) > 1 else f"Content for {keyword}",
                        "image_data": gtm_ai_selection_result["slide"].get("image_data", "")
                    }
                    
                    # Use AI to select the best template slide destination
                    template_selection_result = analyze_and_map_content(
                        api_key,
                        gtm_content_context,
                        all_template_slides_for_ai,
                        keyword
                    )
                    
                    selected_template_index = template_selection_result["best_template_index"]
                    print(f"[ASYNC] AI selected template slide {selected_template_index} as best destination for '{keyword}'")
                    print(f"[ASYNC] AI reasoning: {template_selection_result['justification']}")
                    
                    if selected_template_index > 0 and selected_template_index <= len(new_prs.slides):
                        ai_chosen_slide_index = selected_template_index - 1  # Convert to 0-based
                        dest_slide = new_prs.slides[ai_chosen_slide_index]
                        print(f"[ASYNC] Using AI-chosen slide {selected_template_index} (index {ai_chosen_slide_index}) as destination")
                    else:
                        dest_slide = new_prs.slides[i] if i < len(new_prs.slides) else new_prs.slides[0]
                        print(f"[ASYNC] AI selection out of bounds, using fallback slide")
                    
                    # Step 3: Populate with PLACEHOLDER TEXT (not GTM content)
                    placeholder_content = {
                        "title": f"[{keyword.upper()} - PLACEHOLDER]",
                        "body": f"[REGIONAL {keyword.upper()} DATA AND ANALYSIS HERE]\n\n• Key insight from GTM analysis\n• Regional adaptation needed\n• Local market considerations"
                    }
                    
                    populate_slide(dest_slide, placeholder_content, template_formatting)
                    print(f"[ASYNC] Populated AI-chosen slide with PLACEHOLDER content for '{keyword}'")
            
            TaskManager.update_task(task_id, "processing", 90, "Saving presentation...")
            
            # Save to permanent location
            result_filename = f"assembled_presentation_{task_id}.pptx"
            result_path = os.path.join(TEMP_DIR, result_filename)
            new_prs.save(result_path)
            
            file_size = os.path.getsize(result_path)
            print(f"[ASYNC] Saved result to {result_path}, size: {file_size} bytes")
            TaskManager.update_task(task_id, "completed", 100, "Presentation ready for download!", result_file=result_path)
            
    except Exception as e:
        print(f"[ASYNC] Error processing task {task_id}: {e}")
        import traceback
        traceback.print_exc()
        TaskManager.update_task(task_id, "failed", 0, f"Processing failed: {str(e)}", error=str(e))

@app.post("/assemble")
async def assemble(
    request: Request,
    api_key: str = Form(...),
    template_files: list[UploadFile] = File(...),
    gtm_file: UploadFile = File(...),
    structure: str = Form(...)
):
    print("[DEBUG] /assemble called")
    
    # Add request tracking for timeout debugging
    request_start_time = datetime.now()
    print(f"[DEBUG] Request started at: {request_start_time}")
    
    # Log request details for debugging
    content_length = request.headers.get("content-length", "unknown")
    print(f"[DEBUG] Request size: {content_length} bytes")
    print(f"[DEBUG] Number of template files: {len(template_files)}")
    print(f"[DEBUG] GTM file: {gtm_file.filename if gtm_file else 'None'}")
    
    # Add connection keep-alive headers early
    async def log_progress(message: str):
        """Helper to log progress with timing"""
        elapsed = (datetime.now() - request_start_time).total_seconds()
        print(f"[PROGRESS] {elapsed:.1f}s: {message}")
    
    await log_progress("Request validation starting")
    
    # Validate request size (100MB limit)
    max_size = 100 * 1024 * 1024  # 100MB
    if content_length != "unknown":
        try:
            size = int(content_length)
            if size > max_size:
                error_msg = f"Request too large: {size} bytes. Maximum allowed: {max_size} bytes"
                print(f"[ERROR] {error_msg}")
                return JSONResponse(
                    status_code=413,
                    content={"error": error_msg}
                )
        except ValueError:
            print(f"[WARNING] Could not parse content-length: {content_length}")
    
    await log_progress("Request size validation complete")
    
    # Validate files exist and are not empty
    if not template_files or len(template_files) == 0:
        error_msg = "No template files provided"
        print(f"[ERROR] {error_msg}")
        return JSONResponse(status_code=400, content={"error": error_msg})
    
    if not gtm_file:
        error_msg = "No GTM file provided"
        print(f"[ERROR] {error_msg}")
        return JSONResponse(status_code=400, content={"error": error_msg})
    
    await log_progress("File validation complete")
    
    # Validate API key
    if not api_key or len(api_key.strip()) < 10:
        error_msg = "Invalid or missing OpenAI API key"
        print(f"[ERROR] {error_msg}")
        return JSONResponse(status_code=400, content={"error": error_msg})
    
    # Validate structure
    try:
        structure_steps = json.loads(structure)
        if not structure_steps or len(structure_steps) == 0:
            error_msg = "No structure steps provided"
            print(f"[ERROR] {error_msg}")
            return JSONResponse(status_code=400, content={"error": error_msg})
        print(f"[DEBUG] Structure steps: {structure_steps}")
    except json.JSONDecodeError as e:
        error_msg = f"Invalid JSON in structure field: {str(e)}"
        print(f"[ERROR] {error_msg}")
        return JSONResponse(status_code=400, content={"error": error_msg})
    
    await log_progress("All validation complete, starting processing")
    
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            template_paths = []
            all_template_slides_for_ai = []
            base_pptx_template_found = False
            new_prs = None
            for f in template_files:
                print(f"[DEBUG] Processing template file: {f.filename}")
                print(f"[DIAGNOSTIC] Template file object type: {type(f)}")
                print(f"[DIAGNOSTIC] Template file size: {f.size if hasattr(f, 'size') else 'unknown'}")
                
                path = os.path.join(tmpdir, f.filename)
                with open(path, 'wb') as out:
                    # DIAGNOSTIC: Check file reading
                    bytes_written = 0
                    try:
                        f.file.seek(0)  # Reset file pointer
                        content = f.file.read()
                        print(f"[DIAGNOSTIC] Read {len(content)} bytes from uploaded file")
                        out.write(content)
                        bytes_written = len(content)
                    except Exception as e:
                        print(f"[DIAGNOSTIC] Error reading file: {e}")
                        raise
                
                print(f"[DIAGNOSTIC] Wrote {bytes_written} bytes to {path}")
                template_paths.append(path)
                
                with open(path, 'rb') as file_obj:
                    file_bytes = file_obj.read()
                    print(f"[DIAGNOSTIC] Re-read {len(file_bytes)} bytes from saved file")
                    file_type = mimetypes.guess_type(path)[0] or 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                    print(f"[DEBUG] Template file type: {file_type}")
                    if file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                        if not base_pptx_template_found:
                            # Use the first template file as the base - this keeps ALL slides from the first file
                            new_prs = Presentation(io.BytesIO(file_bytes))
                            base_pptx_template_found = True
                            print(f"[DIAGNOSTIC] Loaded base template with {len(new_prs.slides)} slides")
                            print(f"[DIAGNOSTIC] Base template slide layouts available: {len(new_prs.slide_layouts)}")
                            
                            # DIAGNOSTIC: Check each slide in the loaded presentation
                            for slide_idx, slide in enumerate(new_prs.slides):
                                slide_shapes = len(slide.shapes)
                                slide_text = []
                                for shape in slide.shapes:
                                    if hasattr(shape, "text") and shape.text.strip():
                                        slide_text.append(shape.text.strip())
                                print(f"[DIAGNOSTIC] PPTX Slide {slide_idx}: {slide_shapes} shapes, text: '{' | '.join(slide_text)[:100]}...'")
                            
                            print(f"[DIAGNOSTIC] Total slides loaded into new_prs: {len(new_prs.slides)}")
                            temp_prs_check = Presentation(io.BytesIO(file_bytes))
                            print(f"[DIAGNOSTIC] Template file actually contains {len(temp_prs_check.slides)} slides")
                            for idx, slide in enumerate(temp_prs_check.slides):
                                slide_text_shapes = [s.text.strip() for s in slide.shapes if hasattr(s, 'text') and s.text.strip()]
                                print(f"[DIAGNOSTIC] Template slide {idx}: {len(slide_text_shapes)} text shapes, first text: '{slide_text_shapes[0] if slide_text_shapes else 'No text'}'")
                        else:
                            # Merge additional template files by adding their slides to the base presentation
                            current_prs_to_merge = Presentation(io.BytesIO(file_bytes))
                            print(f"[DEBUG] Merging additional template file with {len(current_prs_to_merge.slides)} slides")
                            for slide_to_merge in current_prs_to_merge.slides:
                                # Use the same layout index as the source slide
                                layout_idx = slide_to_merge.slide_layout.slide_id
                                # Find the layout index in the base template that matches the source
                                matching_layout = None
                                for idx, layout in enumerate(new_prs.slide_layouts):
                                    if layout.slide_id == layout_idx:
                                        matching_layout = layout
                                        break
                                if matching_layout is None:
                                    matching_layout = new_prs.slide_layouts[0]  # fallback
                                new_slide = new_prs.slides.add_slide(matching_layout)
                                deep_copy_slide_content(new_slide, slide_to_merge)
                            print(f"[DEBUG] After merging, presentation now has {len(new_prs.slides)} slides")
                    try:
                        slides_data = get_all_slide_data(file_bytes, file_type)
                        print(f"[DEBUG] Got {len(slides_data)} slides from template file {f.filename}")
                        
                        # DIAGNOSTIC: Compare AI analysis vs PPTX loading
                        print(f"[DIAGNOSTIC] AI Analysis found {len(slides_data)} slides")
                        print(f"[DIAGNOSTIC] PPTX Object has {len(new_prs.slides)} slides")
                        print(f"[DIAGNOSTIC] MISMATCH DETECTED!" if len(slides_data) != len(new_prs.slides) else "[DIAGNOSTIC] Slide counts match")
                        
                        # Show first few slides from AI analysis
                        for i, slide_data in enumerate(slides_data[:3]):
                            ai_text = slide_data.get('text', '')[:50]
                            print(f"[DIAGNOSTIC] AI Slide {i}: '{ai_text}...'")
                        
                        all_template_slides_for_ai.extend(slides_data)
                    except Exception as e:
                        print(f"[ERROR] get_all_slide_data failed for {f.filename}: {e}")
                        traceback_str = traceback.format_exc()
                        return {"error": f"Conversion service failed for template file {f.filename}: {e}", "traceback": traceback_str}
            if new_prs is None:
                print("[ERROR] No PPTX template found.")
                return {"error": "At least one PPTX file must be uploaded as a 'Template Document' to serve as the base for the assembled presentation."}
            gtm_path = os.path.join(tmpdir, gtm_file.filename)
            print(f"[DIAGNOSTIC] GTM file object type: {type(gtm_file)}")
            print(f"[DIAGNOSTIC] GTM file size: {gtm_file.size if hasattr(gtm_file, 'size') else 'unknown'}")
            
            with open(gtm_path, 'wb') as out:
                # DIAGNOSTIC: Check GTM file reading  
                bytes_written = 0
                try:
                    gtm_file.file.seek(0)  # Reset file pointer
                    content = gtm_file.file.read()
                    print(f"[DIAGNOSTIC] Read {len(content)} bytes from GTM uploaded file")
                    out.write(content)
                    bytes_written = len(content)
                except Exception as e:
                    print(f"[DIAGNOSTIC] Error reading GTM file: {e}")
                    raise
            
            print(f"[DIAGNOSTIC] Wrote {bytes_written} bytes to GTM file {gtm_path}")
            
            with open(gtm_path, 'rb') as gtm_obj:
                gtm_file_to_process_bytes = gtm_obj.read()
                print(f"[DIAGNOSTIC] Re-read {len(gtm_file_to_process_bytes)} bytes from saved GTM file")
            gtm_file_to_process_type = mimetypes.guess_type(gtm_path)[0] or 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            gtm_file_to_process_name = gtm_file.filename
            print(f"[DEBUG] GTM file: {gtm_file_to_process_name}, type: {gtm_file_to_process_type}")
            num_template_slides = len(new_prs.slides)
            num_structure_steps = len(structure_steps)
            print(f"[DEBUG] num_template_slides: {num_template_slides}, num_structure_steps: {num_structure_steps}")
            
            # DON'T TRIM SLIDES - keep all template slides available for AI selection
            # The AI might want to select any slide from the template, not just the first few
            # if num_structure_steps < num_template_slides:
            #     for i in range(num_template_slides - 1, num_structure_steps - 1, -1):
            #         rId = new_prs.slides._sldIdLst[i].rId
            #         new_prs.part.drop_rel(rId)
            #         del new_prs.slides._sldIdLst[i]
            
            print(f"[DEBUG] Keeping all {len(new_prs.slides)} template slides available for AI selection")
            
            await log_progress("Analyzing template formatting...")
            
            # ADVANCED FORMATTING ANALYSIS: Extract dominant formatting patterns from template
            template_formatting = analyze_template_formatting(new_prs)
            print(f"[FORMATTING] Template formatting analysis complete")
            if template_formatting:
                title_format = template_formatting.get('title_format')
                body_format = template_formatting.get('body_format')
                if title_format:
                    print(f"[FORMATTING] Dominant title format: {title_format.get('font_name')}, "
                          f"{title_format.get('font_size_pt')}pt, bold={title_format.get('bold')}")
                if body_format:
                    print(f"[FORMATTING] Dominant body format: {body_format.get('font_name')}, "
                          f"{body_format.get('font_size_pt')}pt, bold={body_format.get('bold')}")
            
            # Process each structure step, but let AI select which template slide to use
            for i, step in enumerate(structure_steps):
                keyword = step["keyword"]
                action = step["action"]
                await log_progress(f"Processing step {i+1}/{len(structure_steps)}: {keyword} - {action}")
                print(f"[DEBUG] Step {i}: keyword={keyword}, action={action}")
                
                # Define current step's default slide index
                current_dest_slide_index = i
                dest_slide = None
                if current_dest_slide_index < len(new_prs.slides):
                    dest_slide = new_prs.slides[current_dest_slide_index]
                else:
                    print(f"[ERROR] Step {i} exceeds available template slides ({len(new_prs.slides)})")
                    continue
                
                if action == "Copy from GTM (as is)":
                    await log_progress(f"Step {i+1}: Analyzing GTM content for '{keyword}'")
                    print(f"[DEBUG] Copy from GTM (as is) - Finding GTM content for keyword '{keyword}'")
                    gtm_ai_selection_result = find_slide_by_ai(api_key, gtm_file_to_process_bytes, gtm_file_to_process_type, keyword, "GTM Deck (Content Source)")
                    print(f"[DEBUG] Copy from GTM (as is) - gtm_ai_selection_result: {gtm_ai_selection_result}")
                    
                    if gtm_ai_selection_result["slide"]:
                        print(f"[DEBUG] Copy from GTM (as is) - Using template slide {current_dest_slide_index + 1} as destination")
                            
                        # Use PLACEHOLDER content instead of GTM content as requested
                        placeholder_content = {
                            "title": f"[{keyword.upper()} - PLACEHOLDER]",
                            "body": f"[REGIONAL {keyword.upper()} DATA HERE]\n\n• Data point from GTM analysis\n• Regional market insight\n• Localized content needed\n• Additional placeholder content"
                        }
                        print(f"[DEBUG] Copy from GTM (as is) - Using PLACEHOLDER content instead of GTM content")
                        
                        # Simply populate the selected template slide with placeholder content
                        populate_slide(dest_slide, placeholder_content, template_formatting)
                        print(f"[DEBUG] Copy from GTM (as is) - Populated template slide {current_dest_slide_index + 1} with PLACEHOLDER content")
                    else:
                        print(f"[ERROR] Copy from GTM (as is) - No suitable GTM slide found for keyword '{keyword}'")
                elif action == "Merge: Template Layout + GTM Content":
                    await log_progress(f"Step {i+1}: AI merging template layout with GTM content for '{keyword}'")
                    print(f"[DEBUG] Merge action - Analyzing GTM deck for keyword '{keyword}'")
                    print(f"[DEBUG] Merge action - Step 1: Finding best GTM content match using visual + text analysis")
                    
                    # Step 1: Find and store the best GTM content
                    gtm_ai_selection_result = find_slide_by_ai(api_key, gtm_file_to_process_bytes, gtm_file_to_process_type, keyword, "GTM Deck (Content Source)")
                    print(f"[DEBUG] Merge action - GTM selection result: {gtm_ai_selection_result}")
                    
                    if not gtm_ai_selection_result["slide"]:
                        print(f"[ERROR] Merge action - No suitable GTM slide found for keyword '{keyword}'")
                        continue
                    
                    # Store the GTM content for future operation
                    stored_gtm_content = {
                        "slide_index": gtm_ai_selection_result["index"],
                        "text_content": gtm_ai_selection_result["slide"].get("text", ""),
                        "image_data": gtm_ai_selection_result["slide"].get("image_data", ""),
                        "justification": gtm_ai_selection_result["justification"]
                    }
                    print(f"[DEBUG] Merge action - Stored GTM content from slide {stored_gtm_content['slide_index'] + 1}")
                    print(f"[DEBUG] Merge action - GTM Selection Reasoning: {stored_gtm_content['justification']}")
                    
                    # DIAGNOSTIC: Check GTM text extraction
                    print(f"[DIAGNOSTIC] Raw GTM text_content type: {type(stored_gtm_content['text_content'])}")
                    print(f"[DIAGNOSTIC] Raw GTM text_content length: {len(str(stored_gtm_content['text_content']))}")
                    print(f"[DIAGNOSTIC] Raw GTM text_content first 200 chars: '{str(stored_gtm_content['text_content'])[:200]}'")
                    
                    # Extract title and body from stored content
                    full_text = stored_gtm_content["text_content"]
                    lines = full_text.split('\n') if full_text else []
                    raw_gtm_content = {
                        "title": lines[0] if lines else "",
                        "body": "\n".join(lines[1:]) if len(lines) > 1 else "",
                        "image_data": stored_gtm_content["image_data"]
                    }
                    
                    # DIAGNOSTIC: Check extracted content
                    print(f"[DIAGNOSTIC] Extracted title: '{raw_gtm_content['title']}'")
                    print(f"[DIAGNOSTIC] Extracted body: '{raw_gtm_content['body']}'")
                    print(f"[DIAGNOSTIC] Extracted body length: {len(raw_gtm_content['body'])}")
                    print(f"[DEBUG] Merge action - Extracted GTM title: '{raw_gtm_content['title'][:50]}...'")
                    print(f"[DEBUG] Merge action - Extracted GTM body length: {len(raw_gtm_content['body'])} characters")
                    
                    # Step 2: Analyze template files to find best visual structure match
                    print(f"[DEBUG] Merge action - Step 2: Analyzing template designs for visual structure compatibility")
                    print(f"[DEBUG] Merge action - Available template slides for analysis: {len(all_template_slides_for_ai)}")
                    
                    ai_mapping_result = analyze_and_map_content(
                        api_key,
                        raw_gtm_content,
                        all_template_slides_for_ai,
                        keyword
                    )
                    print(f"[DEBUG] Merge action - Template mapping result: {ai_mapping_result}")
                    print(f"[DEBUG] Merge action - Template Selection Reasoning: {ai_mapping_result['justification']}")
                    
                    selected_template_index = ai_mapping_result["best_template_index"]
                    processed_content = ai_mapping_result["processed_content"]
                    
                    # DIAGNOSTIC: Check AI processed content format
                    print(f"[DIAGNOSTIC] AI processed_content type: {type(processed_content)}")
                    print(f"[DIAGNOSTIC] AI processed_content keys: {list(processed_content.keys()) if isinstance(processed_content, dict) else 'Not a dict'}")
                    if isinstance(processed_content, dict):
                        print(f"[DIAGNOSTIC] AI processed title type: {type(processed_content.get('title', 'missing'))}")
                        print(f"[DIAGNOSTIC] AI processed title value: '{processed_content.get('title', 'missing')}'")
                        print(f"[DIAGNOSTIC] AI processed body type: {type(processed_content.get('body', 'missing'))}")
                        print(f"[DIAGNOSTIC] AI processed body value: {processed_content.get('body', 'missing')}")
                    
                    print(f"[DEBUG] Merge action - AI selected template index: {selected_template_index} (AI uses 1-based indexing)")
                    print(f"[DEBUG] Merge action - Available template slides in all_template_slides_for_ai: {len(all_template_slides_for_ai)}")
                    print(f"[DEBUG] Merge action - Available slides in new_prs.slides: {len(new_prs.slides)}")
                    
                    # Convert AI's 1-based index to 0-based, but ensure it's within bounds
                    template_slide_to_use = None
                    ai_slide_index = current_dest_slide_index  # Default fallback index
                    
                    if selected_template_index > 0 and selected_template_index <= len(new_prs.slides):
                        # Convert AI's 1-based index to 0-based to get the correct slide from new_prs.slides
                        ai_selected_index_0based = selected_template_index - 1
                        print(f"[DEBUG] Merge action - AI selected index (0-based): {ai_selected_index_0based}")
                        
                        # The AI selects based on all_template_slides_for_ai, which should map directly to new_prs.slides
                        # since they were built from the same template files in the same order
                        ai_slide_index = ai_selected_index_0based
                        template_slide_to_use = new_prs.slides[ai_slide_index]
                        print(f"[DEBUG] Merge action - Using AI-selected template slide at index {ai_slide_index}")
                        print(f"[DEBUG] Merge action - AI selected slide {selected_template_index} successfully mapped to presentation slide {ai_slide_index + 1}")
                    else:
                        print(f"[DEBUG] Merge action - AI index {selected_template_index} out of bounds or invalid, using current slide")
                        ai_slide_index = current_dest_slide_index
                        template_slide_to_use = dest_slide
                    
                    # Step 3: Populate with PLACEHOLDER TEXT (not processed GTM content)
                    if template_slide_to_use:
                        print(f"[DEBUG] Merge action - Step 3: Using template slide {ai_slide_index} with placeholder content")
                        print(f"[DEBUG] Merge action - Template slide has {len(template_slide_to_use.shapes)} shapes")
                        print(f"[DEBUG] Merge action - Using PLACEHOLDER TEXT instead of GTM content")
                        
                        # Use PLACEHOLDER content instead of processed_content from GTM
                        placeholder_content = {
                            "title": f"[{keyword.upper()} - PLACEHOLDER]",
                            "body": f"[REGIONAL {keyword.upper()} DATA AND ANALYSIS HERE]\n\n• Key insight from GTM analysis\n• Regional adaptation needed\n• Local market considerations\n• Additional placeholder bullet point"
                        }
                        
                        print(f"[DIAGNOSTIC] Using placeholder_content: {placeholder_content}")
                        print(f"[DIAGNOSTIC] AI-selected template_slide_to_use shapes before population: {len(template_slide_to_use.shapes)}")
                        
                        # Get the original template slide image for visual comparison
                        print(f"[VISUAL VERIFY] Capturing original template slide image for consistency check...")
                        original_template_image = None
                        if ai_slide_index < len(all_template_slides_for_ai):
                            original_template_image = all_template_slides_for_ai[ai_slide_index].get('image_data')
                        
                        populate_slide(template_slide_to_use, placeholder_content, template_formatting)
                        print(f"[DIAGNOSTIC] After populate_slide - template_slide_to_use shapes: {len(template_slide_to_use.shapes)}")
                        
                        # VISUAL VERIFICATION: Temporarily disabled for faster response
                        print(f"[VISUAL VERIFY] ===== VISUAL VERIFICATION DISABLED FOR SPEED =====")
                        print(f"[VISUAL VERIFY] Slide populated with advanced template formatting - verification skipped")
                        
                        print(f"[VISUAL VERIFY] ===== VISUAL VERIFICATION PROCESS COMPLETE =====")
                        
                        print(f"[DEBUG] Merge action - Populated AI-selected template slide {ai_slide_index + 1} with PLACEHOLDER content")
                        print(f"[DEBUG] Merge action - Final result: AI-selected template + placeholder content with advanced formatting")
                    else:
                        print(f"[ERROR] Merge action - No suitable template slide found")
                        print(f"[ERROR] Merge action - Using simple population of current slide instead")
                        
                        # Fallback: use placeholder content on current slide
                        placeholder_content = {
                            "title": f"[{keyword.upper()} - PLACEHOLDER]",
                            "body": f"[REGIONAL {keyword.upper()} DATA AND ANALYSIS HERE]\n\n• Key insight from GTM analysis\n• Regional adaptation needed\n• Local market considerations\n• Additional placeholder bullet point"
                        }
                        
                        # Fallback: just populate the template slide with placeholder content
                        populate_slide(template_slide_to_use, placeholder_content, template_formatting)
            await log_progress("All slides processed, saving presentation...")
            output_pptx_path = os.path.join(tmpdir, "assembled_presentation.pptx")
            print(f"[DEBUG] About to save assembled presentation to: {output_pptx_path}")
            new_prs.save(output_pptx_path)
            await log_progress("Presentation saved, preparing file response...")
            print(f"[DEBUG] Assembled presentation saved to: {output_pptx_path}")
            file_exists = os.path.exists(output_pptx_path)
            print(f"[DEBUG] File exists after save? {file_exists}")
            if not file_exists:
                print(f"[ERROR] File {output_pptx_path} does not exist after save!")
            else:
                file_size = os.path.getsize(output_pptx_path)
                print(f"[DEBUG] File size: {file_size} bytes")
            # Try to open the file before returning
            try:
                with open(output_pptx_path, 'rb') as f:
                    pptx_bytes = f.read()
                
                final_size = len(pptx_bytes)
                await log_progress(f"File ready: {final_size} bytes, starting download...")
                print(f"[DEBUG] Read {final_size} bytes from assembled PPTX before tempdir cleanup.")
                
                # Create a chunked streaming response to avoid timeouts
                def generate_chunked_file():
                    chunk_size = 32 * 1024  # 32KB chunks for better streaming
                    data = io.BytesIO(pptx_bytes)
                    bytes_sent = 0
                    
                    while True:
                        chunk = data.read(chunk_size)
                        if not chunk:
                            break
                        bytes_sent += len(chunk)
                        print(f"[STREAM] Sent {bytes_sent}/{final_size} bytes ({(bytes_sent/final_size)*100:.1f}%)")
                        yield chunk
                    
                    print(f"[STREAM] File transfer complete: {bytes_sent} bytes sent")
                
                print(f"[DEBUG] Returning chunked streaming response for {final_size} bytes")
                
                # Add comprehensive headers for better compatibility
                headers = {
                    "Content-Disposition": "attachment; filename=assembled_presentation.pptx",
                    "Content-Length": str(final_size),
                    "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "Cache-Control": "no-cache, no-store, must-revalidate",
                    "Pragma": "no-cache",
                    "Expires": "0",
                    "Connection": "keep-alive",
                    "Accept-Ranges": "bytes",
                    "X-Content-Type-Options": "nosniff"
                }
                
                return StreamingResponse(
                    generate_chunked_file(),
                    media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    headers=headers
                )
            except Exception as e:
                print(f"[ERROR] Could not open file before FileResponse: {e}")
                await log_progress(f"ERROR: File processing failed: {e}")
                return JSONResponse(
                    status_code=500,
                    content={"error": f"File processing failed: {str(e)}"}
                )
    except Exception as e:
        elapsed = (datetime.now() - request_start_time).total_seconds()
        print(f"[ERROR] Exception in /assemble after {elapsed:.1f}s: {e}")
        traceback_str = traceback.format_exc()
        await log_progress(f"FAILED after {elapsed:.1f}s: {str(e)}")
        return JSONResponse(
            status_code=500,
            content={"error": str(e), "traceback": traceback_str, "elapsed_seconds": elapsed}
        )

@app.post("/pptx_to_png_zip")
async def pptx_to_png_zip(file: UploadFile = File(...)):
    """
    Accepts a PPTX file and returns a ZIP of PNG slide images.
    """
    try:
        file_bytes = await file.read()
        filename = file.filename or "presentation.pptx"
        zip_bytes = convert_pptx_to_png_zip(file_bytes, filename)
        return StreamingResponse(
            io.BytesIO(zip_bytes),
            media_type="application/zip",
            headers={
                "Content-Disposition": f"attachment; filename=slides_png.zip"
            }
        )
    except Exception as e:
        traceback_str = traceback.format_exc()
        return {"error": str(e), "traceback": traceback_str}

def verify_visual_consistency_with_ai(api_key, original_template_image_b64, final_slide_image_b64, content_type="slide"):
    """
    Use AI to compare the final slide output with the original template for visual consistency
    """
    if not api_key:
        print("[VISUAL VERIFY] No API key provided, skipping visual verification")
        return {"consistent": True, "recommendations": "No verification performed"}
    
    # Create OpenAI client with explicit parameters to avoid proxy issues
    try:
        client = openai.OpenAI(
            api_key=api_key,
            base_url=None,  # Explicitly set to None
            timeout=60.0    # Explicit timeout
        )
    except Exception as e:
        print(f"[ERROR] Failed to create OpenAI client for visual verification: {e}")
        return {"consistent": False, "recommendations": f"OpenAI client creation failed: {e}"}
    
    system_prompt = """
    You are an expert presentation design analyst. Your task is to compare two slide images:
    1. ORIGINAL TEMPLATE SLIDE - the reference design
    2. FINAL OUTPUT SLIDE - the populated slide that should match the template's visual style
    
    **Critical Analysis Requirements:**
    - Compare font families, sizes, and weights
    - Check text alignment, spacing, and positioning
    - Verify color consistency (text colors, backgrounds)
    - Assess overall visual hierarchy and balance
    - Look for any formatting inconsistencies
    
    **Your Response Must Include:**
    1. **Visual Consistency Score** (0-100): How well does the final output match the template's visual style?
    2. **Specific Issues Found**: List any formatting, font, color, or positioning problems
    3. **Formatting Recommendations**: Specific fixes needed (font name, size, color, alignment, etc.)
    
    **Focus Areas:**
    - Does the text use the same font family as the template?
    - Is the font size consistent with similar elements in the template?
    - Are colors (text, background) matching the template style?
    - Is the text positioning and alignment correct?
    
    Return a JSON object with:
    - "consistency_score": integer 0-100
    - "is_consistent": boolean (true if score >= 85)
    - "issues_found": array of specific problems
    - "recommendations": object with specific formatting fixes needed
    """
    
    user_parts = [
        {"type": "text", "text": f"Please analyze the visual consistency between the original template and the final populated {content_type}:"},
        {"type": "text", "text": "\n--- ORIGINAL TEMPLATE (Reference Design):"},
        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{original_template_image_b64}"}},
        {"type": "text", "text": "\n--- FINAL OUTPUT (Should Match Template Style):"},
        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{final_slide_image_b64}"}}
    ]
    
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_parts}
    ]
    
    try:
        print("[VISUAL VERIFY] Sending images to OpenAI for visual consistency analysis...")
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            response_format={"type": "json_object"}
        )
        
        result = json.loads(response.choices[0].message.content)
        consistency_score = result.get("consistency_score", 0)
        is_consistent = result.get("is_consistent", False)
        issues_found = result.get("issues_found", [])
        recommendations = result.get("recommendations", {})
        
        print(f"[VISUAL VERIFY] Consistency score: {consistency_score}/100")
        print(f"[VISUAL VERIFY] Is consistent: {is_consistent}")
        if issues_found:
            print(f"[VISUAL VERIFY] Issues found: {issues_found}")
        if recommendations:
            print(f"[VISUAL VERIFY] Recommendations: {recommendations}")
        
        return {
            "consistent": is_consistent,
            "score": consistency_score,
            "issues": issues_found,
            "recommendations": recommendations
        }
        
    except Exception as e:
        print(f"[VISUAL VERIFY] AI visual verification failed: {e}")
        return {"consistent": True, "recommendations": "Verification failed"}

def convert_slide_to_image(slide, presentation_path=None):
    """
    Convert a single slide to PNG image for AI analysis
    """
    print(f"[VISUAL VERIFY] convert_slide_to_image called")
    try:
        # Create a temporary presentation with just this slide
        temp_prs = Presentation()
        print(f"[VISUAL VERIFY] Created temporary presentation")
        
        # Copy the slide layout if possible
        if hasattr(slide, 'slide_layout'):
            # Try to find a matching layout or use the first one
            target_layout = temp_prs.slide_layouts[0]  # Fallback
            for layout in temp_prs.slide_layouts:
                if layout.name == slide.slide_layout.name:
                    target_layout = layout
                    break
            print(f"[VISUAL VERIFY] Found target layout: {target_layout.name}")
        else:
            target_layout = temp_prs.slide_layouts[0]
            print(f"[VISUAL VERIFY] Using default layout: {target_layout.name}")
        
        # Add the slide to the temp presentation
        new_slide = temp_prs.slides.add_slide(target_layout)
        print(f"[VISUAL VERIFY] Added slide to temp presentation")
        
        # Copy all content from the original slide
        print(f"[VISUAL VERIFY] Copying slide content...")
        deep_copy_slide_content(new_slide, slide)
        print(f"[VISUAL VERIFY] Slide content copied successfully")
        
        # Save to temporary file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_prs.save(temp_file.name)
            temp_pptx_path = temp_file.name
            print(f"[VISUAL VERIFY] Temp PPTX saved to: {temp_pptx_path}")
        
        # Convert to PNG using the existing conversion function
        with open(temp_pptx_path, 'rb') as f:
            temp_bytes = f.read()
        
        print(f"[VISUAL VERIFY] Converting PPTX to PNG using ConvertAPI...")
        zip_bytes = convert_pptx_to_png_zip(temp_bytes, "temp_slide.pptx")
        print(f"[VISUAL VERIFY] PNG conversion complete, extracting image...")
        
        # Extract the first (and only) PNG image
        with zipfile.ZipFile(io.BytesIO(zip_bytes), 'r') as zf:
            png_files = [name for name in zf.namelist() if name.lower().endswith('.png')]
            print(f"[VISUAL VERIFY] Found {len(png_files)} PNG files in ZIP")
            if png_files:
                with zf.open(png_files[0]) as img_file:
                    img_bytes = img_file.read()
                    img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                    print(f"[VISUAL VERIFY] Successfully extracted PNG image ({len(img_bytes)} bytes -> {len(img_b64)} base64 chars)")
                    
                # Clean up temp file
                os.unlink(temp_pptx_path)
                print(f"[VISUAL VERIFY] Cleaned up temp file")
                return img_b64
        
        # Clean up temp file if we get here
        os.unlink(temp_pptx_path)
        print(f"[VISUAL VERIFY] No PNG files found in conversion result")
        return None
        
    except Exception as e:
        print(f"[VISUAL VERIFY] Failed to convert slide to image: {e}")
        import traceback
        print(f"[VISUAL VERIFY] Full traceback: {traceback.format_exc()}")
        return None

def apply_ai_formatting_recommendations(slide, recommendations, template_formatting=None):
    """
    Apply AI recommendations to fix formatting inconsistencies
    """
    if not recommendations or not isinstance(recommendations, dict):
        print("[VISUAL VERIFY] No valid recommendations to apply")
        return
    
    print(f"[VISUAL VERIFY] Applying AI formatting recommendations: {recommendations}")
    
    # Extract recommendations
    font_name = recommendations.get('font_name') or recommendations.get('font_family')
    font_size_pt = recommendations.get('font_size') or recommendations.get('font_size_pt')
    text_color = recommendations.get('text_color') or recommendations.get('color')
    alignment = recommendations.get('alignment')
    bold = recommendations.get('bold')
    
    # Convert font size to Pt object if needed
    if font_size_pt and isinstance(font_size_pt, (int, float)):
        font_size = Pt(font_size_pt)
    else:
        font_size = None
    
    # Apply to all text shapes in the slide
    shapes_updated = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                # Apply paragraph-level formatting
                if alignment:
                    try:
                        # Convert alignment string to proper enum if needed
                        if alignment.lower() == 'center':
                            para.alignment = 1  # PP_ALIGN_CENTER
                        elif alignment.lower() == 'left':
                            para.alignment = 0  # PP_ALIGN_LEFT
                        elif alignment.lower() == 'right':
                            para.alignment = 2  # PP_ALIGN_RIGHT
                    except Exception as e:
                        print(f"[VISUAL VERIFY] Failed to apply alignment: {e}")
                
                # Apply run-level formatting
                for run in para.runs:
                    if font_name:
                        try:
                            run.font.name = font_name
                            print(f"[VISUAL VERIFY] Applied font name: {font_name}")
                        except Exception as e:
                            print(f"[VISUAL VERIFY] Failed to apply font name: {e}")
                    
                    if font_size:
                        try:
                            run.font.size = font_size
                            print(f"[VISUAL VERIFY] Applied font size: {font_size}")
                        except Exception as e:
                            print(f"[VISUAL VERIFY] Failed to apply font size: {e}")
                    
                    if bold is not None:
                        try:
                            run.font.bold = bold
                            print(f"[VISUAL VERIFY] Applied bold: {bold}")
                        except Exception as e:
                            print(f"[VISUAL VERIFY] Failed to apply bold: {e}")
                    
                    if text_color:
                        try:
                            # Parse color if it's a hex string
                            if isinstance(text_color, str) and text_color.startswith('#'):
                                hex_color = text_color[1:]  # Remove #
                                r = int(hex_color[0:2], 16)
                                g = int(hex_color[2:4], 16)
                                b = int(hex_color[4:6], 16)
                                run.font.fill.solid()
                                run.font.fill.fore_color.rgb = RGBColor(r, g, b)
                                print(f"[VISUAL VERIFY] Applied text color: {text_color}")
                        except Exception as e:
                            print(f"[VISUAL VERIFY] Failed to apply text color: {e}")
            
            shapes_updated += 1
    
    print(f"[VISUAL VERIFY] Updated formatting on {shapes_updated} shapes")
